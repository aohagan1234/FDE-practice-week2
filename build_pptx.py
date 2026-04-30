"""Generate ATX Assessment Summary presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
ORANGE = RGBColor(0xE3, 0x6B, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
MID    = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xF2, 0xF4, 0xF7)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
RED    = RGBColor(0xBF, 0x35, 0x2A)
AMBER  = RGBColor(0xCC, 0x84, 0x00)
TEAL   = RGBColor(0x0E, 0x7C, 0x7B)
LBLUE  = RGBColor(0xBB, 0xCC, 0xEE)
DBLUE2 = RGBColor(0x23, 0x52, 0x96)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(s, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def mtxt(s, x, y, w, h, lines):
    """Multi-line textbox. lines = list of (text, size, bold, color, align, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align, italic) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.italic = italic
    return tb


def header(s, title, subtitle=None):
    rect(s, 0, 0, 13.33, 1.28, NAVY)
    txt(s, title, 0.38, 0.07, 12.5, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        txt(s, subtitle, 0.38, 0.77, 12.5, 0.44, size=12, color=LBLUE)


def slide_number(s, n):
    txt(s, str(n), 12.9, 7.15, 0.35, 0.3, size=9, color=MID, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 4.95, 13.33, 2.55, DBLUE2)
rect(s, 0, 2.88, 13.33, 0.07, ORANGE)   # orange divider

txt(s, "ALDRIDGE & SYKES  ·  HR OPERATIONS  ·  WEEK 2 ATX ASSESSMENT",
    0.55, 0.35, 12.2, 0.45, size=10, color=LBLUE, bold=True)
txt(s, "Should we build an AI agent to replace\nmanual onboarding coordination?",
    0.55, 0.82, 12.0, 1.95, size=38, bold=True, color=WHITE)
txt(s, "A structured 7-phase analysis: what to automate, what not to, and what was built.",
    0.55, 3.0, 11.5, 0.55, size=14, color=LBLUE)

# Answer callout
rect(s, 0.55, 5.1, 5.2, 1.6, ORANGE)
txt(s, "The conclusion:", 0.72, 5.15, 4.8, 0.38, size=11, bold=True, color=WHITE)
txt(s, "Yes — automate 50% fully.\nOversee 35%. Keep 15% human-only.",
    0.72, 5.48, 4.8, 1.1, size=16, bold=True, color=WHITE)

# Key stats
mtxt(s, 6.0, 5.42, 6.9, 1.7, [
    ("220+ hires / year", 13, True,  LBLUE,  PP_ALIGN.LEFT, False),
    ("40 tasks per hire across 5 disconnected systems", 12, False, LBLUE, PP_ALIGN.LEFT, False),
    ("550 hours / year spent on manual coordination", 12, False, LBLUE, PP_ALIGN.LEFT, False),
    ("$2,507 federal fine per missed I-9 deadline", 12, False, RGBColor(0xFF,0xAA,0x88), PP_ALIGN.LEFT, False),
])


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "The Problem", "What HR Ops is actually dealing with — in plain numbers")
slide_number(s, 2)

# Stat boxes (left column)
stat_data = [
    (NAVY,   "550 hours / year",  "One person's entire working year\nspent purely on coordination"),
    (NAVY,   "40 tasks per hire", "Across 5 systems that don't\ntalk to each other"),
    (RED,    "$2,507 fine",       "Per missed I-9 deadline (federal law)\n2 violations have already occurred"),
    (ORANGE, "3 people",          "Acting as the human glue between\n5 disconnected systems"),
]
for i, (color, title, body) in enumerate(stat_data):
    row, col = i // 2, i % 2
    bx = 0.35 + col * 2.15
    by = 1.42 + row * 1.62
    rect(s, bx, by, 1.98, 1.48, color)
    txt(s, title, bx+0.1, by+0.08, 1.78, 0.5,  size=15, bold=True, color=WHITE)
    txt(s, body,  bx+0.1, by+0.58, 1.78, 0.82, size=10, color=WHITE)

# Right: narrative
mtxt(s, 4.65, 1.38, 8.3, 5.7, [
    ("Every time someone joins Aldridge & Sykes:", 15, True,  NAVY, PP_ALIGN.LEFT, False),
    ("",                                            6,  False, DARK, PP_ALIGN.LEFT, False),
    ("  •  HR Ops manually checks 40 tasks across 5 systems",   12, False, DARK, PP_ALIGN.LEFT, False),
    ("  •  None of the systems share data automatically",        12, False, DARK, PP_ALIGN.LEFT, False),
    ("  •  Someone has to open each system and check individually", 12, False, DARK, PP_ALIGN.LEFT, False),
    ("  •  Overdue tasks are chased by email — manually",        12, False, DARK, PP_ALIGN.LEFT, False),
    ("  •  I-9 deadlines are tracked on a phone Notes app",      12, False, DARK, PP_ALIGN.LEFT, False),
    ("",                                            6,  False, DARK, PP_ALIGN.LEFT, False),
    ("That is 2.5 hours per hire × 220 hires = 550 hours of pure\ncoordination work per year.", 13, True, NAVY, PP_ALIGN.LEFT, False),
    ("",                                            6,  False, DARK, PP_ALIGN.LEFT, False),
    ("Roughly a quarter of each person's working year. Spent on\ntasks a well-designed system could do automatically.", 12, False, DARK, PP_ALIGN.LEFT, True),
    ("",                                            6,  False, DARK, PP_ALIGN.LEFT, False),
    ("The I-9 compliance risk is the most urgent issue. Missing\nthe 3-day federal deadline costs $2,507 per hire. Two\nviolations have already occurred under the current process.", 12, True, RED, PP_ALIGN.LEFT, False),
])


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE METHOD
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "The Method — 7 Phases", "From 'what don't we know?' to 'does it work in code?'")
slide_number(s, 3)

txt(s, "Rather than jumping straight to building, this used the ATX methodology: a 7-phase framework for deciding where an AI agent adds genuine value — and where it doesn't.",
    0.38, 1.33, 12.55, 0.55, size=12, color=DARK, italic=True)

phases = [
    (NAVY,  "1  Discovery",            "Asked 6 design-changing questions about the real process.\nFound that Priya's Excel tracker — not Workday — is the true\noperational source of truth."),
    (NAVY,  "2  Cognitive Load Map",   "Decomposed all 40 tasks into decisions vs. rule-following.\nIdentified 8 control handoff points — where the agent must\nstop and pass control back to a human."),
    (NAVY,  "3  Delegation Scoring",   "Scored all 9 task clusters across 5 dimensions: input\nstructure, decision predictability, exception rate, tool\ncoverage, and risk if wrong."),
    (BLUE,  "4  Volume × Value",       "Ranked which tasks are worth automating first. Task\nmonitoring scores highest: 3,500 events/year, mostly\nrule-based. I-9 monitoring is non-negotiable regardless."),
    (BLUE,  "5  Agent Specification",  "Wrote the full agent spec: 9 activities, decision rules,\nescalation workflows, KPIs, and failure modes. Revised\nafter the build loop found 6 gaps in the original draft."),
    (TEAL,  "6  Build Loop",           "Fed the spec to Claude, asked what could be built confidently\nvs. what needed clarification. Built all confident parts in\nPython. 37 automated tests, all passing."),
    (TEAL,  "7  System Inventory",     "Documented every integration in detail. Key finding: Saba\nLMS has no API. Weekly batch file only. This changes the\ndetection window from 2 hours to up to 7 days."),
]

positions = [
    (0.35, 1.98), (4.55, 1.98), (8.75, 1.98),
    (0.35, 3.85), (4.55, 3.85), (8.75, 3.85),
    (0.35, 5.72),
]
for (bx, by), (color, title, body) in zip(positions, phases):
    bw = 3.9 if bx < 8.0 else 3.9
    rect(s, bx, by, bw, 0.44, color)
    txt(s, title, bx+0.12, by+0.04, bw-0.2, 0.36, size=12, bold=True, color=WHITE)
    rect(s, bx, by+0.44, bw, 1.25, WHITE)
    txt(s, body,  bx+0.12, by+0.52, bw-0.2, 1.14, size=10, color=DARK)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT WE FOUND
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "What We Found", "Not everything should be automated — and that split is intentional")
slide_number(s, 4)

cols = [
    (GREEN,  "~50%  Fully Automated",
     "Agent acts alone. No human review\nrequired per action.\n\n"
     "•  Task monitoring & reminders\n"
     "•  I-9 compliance escalation\n"
     "•  Manager handoff notification\n"
     "•  Deadline calculation\n\n"
     "These tasks follow clear rules, are low-risk if\nwrong, and happen hundreds of times a year."),
    (AMBER,  "~35%  Agent Proposes,\nHuman Approves",
     "Agent generates a recommendation.\nPriya reviews and confirms.\n\n"
     "•  Compliance training track\n"
     "•  Buddy matching (ranked list)\n"
     "•  IT access package selection\n\n"
     "Designed for Phase 2. Not in the\ncurrent build. Requires a structured\napproval workflow to be set up first."),
    (RED,    "~15%  Human Only",
     "No agent involvement.\nToo consequential or irreversible.\n\n"
     "•  Hold decisions\n"
     "•  Hire type classification\n"
     "•  Escalation responses\n\n"
     "Agent detects the trigger and alerts.\nHuman decides what happens next.\nThis boundary is non-negotiable."),
]

for i, (color, title, body) in enumerate(cols):
    cx = 0.38 + i * 4.32
    rect(s, cx, 1.38, 4.1,  0.8,  color)
    txt(s,  title, cx+0.15, 1.43, 3.82, 0.72, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, cx, 2.18, 4.1,  3.92, WHITE)
    txt(s,  body,  cx+0.15, 2.28, 3.82, 3.72, size=11, color=DARK)

rect(s, 0.38, 6.25, 12.55, 1.0, NAVY)
txt(s, "ROI is honest:  Financial payback is 18–24 months (only 73 hires per person per year). "
       "The real justification is compliance risk: one prevented I-9 violation ($2,507 penalty) "
       "covers the entire build cost. The agent is justified — but the analysis doesn't pretend the numbers are stronger than they are.",
    0.55, 6.32, 12.25, 0.87, size=11, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE 7 DELIVERABLES
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "The 7 Deliverables", "One output per phase — each one builds on the last")
slide_number(s, 5)

deliverables = [
    (NAVY,   "D1", "Discovery Questions",
     "6 questions that probe the lived work — not just the SOP. Simulated COO answers revealed Priya's Excel tracker (the real source of truth), a Legacy HR system not in the official stack, and an iPhone Notes app used for I-9 monitoring."),
    (NAVY,   "D2", "Cognitive Load Map",
     "Every HR Ops task broken down into decisions vs. rule-following, with the data source, output, and pause points for each. Identified 8 control handoff points where the agent must hand back to a human."),
    (NAVY,   "D3", "Delegation Suitability Matrix",
     "9 task clusters scored across 5 dimensions. Shows which clusters an agent can own and which must stay human — including hold decisions (scored 2.0/5.0 — Human Only) and I-9 monitoring (4.8/5.0 — Fully Agentic, mandatory)."),
    (BLUE,   "D4", "Volume × Value Analysis",
     "Ranked clusters by how often they occur vs. how much an agent can reduce the effort. Task monitoring scores highest (3,500 events/year). No cluster hits the top tier individually — Agent 1 is a bundle of the best clusters."),
    (TEAL,   "D5", "Agent Purpose Document",
     "Full buildable specification: 9 activities, autonomy matrix, escalation workflows, KPIs, failure modes, and data schemas. Revised after the build loop found 6 gaps — including a critical error: I-9 was set to poll daily, not every 2 hours."),
    (TEAL,   "D6", "System & Data Inventory",
     "Every integration documented with API type, rate limits, data quality notes, and risks. Critical finding: Saba LMS has no API — weekly SFTP batch only. Corrects the spec's original SLA claim of '4 hours' for compliance training detection."),
    (GREEN,  "Build", "Python Implementation + 37 Tests",
     "Working code: 5 Python files implementing all 9 activities. Idempotent escalations (no alert flooding), rate-limited reminders (max 1 per 24h per task), I-9 monitoring every 2 hours. 37 automated tests, all passing. Run: pytest tests/ -v"),
]

row_h = 0.745
for i, (color, label, title, body) in enumerate(deliverables):
    by = 1.35 + i * row_h
    rect(s, 0.28, by, 0.72, 0.63, color)
    txt(s, label, 0.28, by+0.1, 0.72, 0.43, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, 1.1,  by+0.04, 3.15, 0.38, size=12, bold=True, color=NAVY)
    txt(s, body,  4.35, by+0.02, 8.6,  0.65, size=10, color=DARK)
    if i < 6:
        rect(s, 0.28, by+0.68, 12.77, 0.015, RGBColor(0xCC, 0xCC, 0xCC))


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WHAT WAS BUILT
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "What Was Built — Agent 1: Coordination Orchestrator",
       "One agent, always on, polling every 2 hours across all 5 systems")
slide_number(s, 6)

# Left: what it does
rect(s, 0.35, 1.38, 5.95, 0.46, GREEN)
txt(s, "What the agent does automatically", 0.5, 1.41, 5.72, 0.38, size=12, bold=True, color=WHITE)

does_items = [
    "Reads all 5 systems every 2 hours for current task status",
    "Flags tasks overdue by more than 24 hours",
    "Sends email reminders to task owners — max 1 per task per 24 hours",
    "Escalates I-9 risk on Day 2 (HIGH) and Day 3 (CRITICAL) — no exceptions",
    "Submits IT provisioning requests when hire's role is in the access matrix",
    "Notifies the hiring manager when Day 1 readiness is confirmed",
    "Routes escalations by severity: email / SMS / phone call",
    "Logs every action to a structured audit trail",
]
for j, item in enumerate(does_items):
    by = 1.92 + j * 0.53
    rect(s, 0.35, by+0.04, 0.28, 0.28, GREEN)
    txt(s, "✓", 0.35, by+0.04, 0.28, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, item, 0.72, by+0.03, 5.48, 0.44, size=10.5, color=DARK)

# Right: what it does NOT do
rect(s, 6.55, 1.38, 6.42, 0.46, RED)
txt(s, "What remains human-only", 6.7, 1.41, 6.18, 0.38, size=12, bold=True, color=WHITE)

not_items = [
    ("Hold decisions",           "Irreversible employment action with legal implications"),
    ("Hire type classification", "5–10% of cases are genuinely ambiguous; cascades if wrong"),
    ("Final buddy selection",    "Team dynamics and interpersonal fit can't be codified"),
    ("Compliance track approval","Policy interpretation needed on low-confidence matches"),
    ("Escalation responses",     "Agent alerts; human decides the action and owns the outcome"),
]
for j, (title, reason) in enumerate(not_items):
    by = 1.92 + j * 0.8
    rect(s, 6.55, by+0.04, 0.28, 0.28, RED)
    txt(s, "✗", 6.55, by+0.04, 0.28, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title,  6.92, by+0.02, 5.9, 0.32, size=11, bold=True,  color=DARK)
    txt(s, reason, 6.92, by+0.3,  5.9, 0.42, size=10, italic=True, color=MID)

# Code bar
rect(s, 0.35, 6.47, 12.6, 0.78, NAVY)
txt(s, "Code:  5 Python files  ·  ~650 lines  ·  37 automated tests, all passing  ·  python main.py to run  ·  pytest tests/ -v to verify",
    0.55, 6.64, 12.3, 0.44, size=11.5, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — KEY DESIGN DECISIONS
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "Key Design Decisions", "The non-obvious choices — and why they were made this way")
slide_number(s, 7)

decisions = [
    (ORANGE, "I-9 monitoring is fully automated despite being CRITICAL risk",
     "Counter-intuitive. The highest-risk item is the most appropriate for automation. I-9 escalation is deterministic (Day 2, Day 3 — no judgment), non-deferrable (federal law), and exactly where human memory fails under load. An agent checking every 2 hours is strictly better than a person who might forget during a busy week with 5 starters."),
    (RED,    "Hold decisions are human-only — always, with no exceptions",
     "Hold = irreversible employment action. Legal implications. The agent detects the trigger and fires an escalation. It does not write hire.status = ON_HOLD. This boundary is non-negotiable regardless of how confident the signal looks. The COO confirmed hold decisions currently happen on WhatsApp — the Workday record is retrospective."),
    (TEAL,   "Saba LMS has no API — and the spec had to be corrected",
     "The original spec described Saba LMS as a 'SOAP API (legacy).' The enriched scenario explicitly says 'no API.' The only path is a weekly SFTP batch file. This changes compliance training detection from 2 hours to up to 7 days. Rather than hide this, the spec was corrected and a three-state batch handling model was built into the design."),
    (BLUE,   "The ROI case is honest rather than flattering",
     "Volume is genuinely low: 73 hires per person per year. Year 1 return on investment is negative. This was kept in the analysis rather than corrected to look better. The real case for building is compliance de-risking (one prevented I-9 violation at $2,507 covers the build cost) and scalability — not pure time savings. Agent 2 scores higher in the formula but depends on Agent 1 first."),
]

for i, (color, title, body) in enumerate(decisions):
    row, col = i // 2, i % 2
    bx = 0.35 + col * 6.52
    by = 1.38 + row * 2.75
    rect(s, bx, by, 6.12, 0.46, color)
    txt(s, title, bx+0.12, by+0.05, 5.9, 0.36, size=11, bold=True, color=WHITE)
    rect(s, bx, by+0.46, 6.12, 2.08, WHITE)
    txt(s, body,  bx+0.12, by+0.56, 5.9,  1.9,  size=10.5, color=DARK)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BEFORE GOING LIVE
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "Before Going Live", "Three things that must happen before the agent is deployed")
slide_number(s, 8)

steps = [
    (RED,    "1", "Get Priya's Excel tracker and treat it as the design baseline",
     "The tracker is the real system of record — not Workday. Workday is the formal record; the tracker is the live operational view, updated daily. Every column it has that Workday doesn't is a gap the agent must fill. Any design that ignores it will produce an agent with less situational awareness than Priya currently has."),
    (AMBER,  "2", "Fix the escalation channel — Teams or SMS, not email",
     "Critical alerts sent by email won't be seen for 3–4 hours. Urgent decisions happen on Teams or by phone. The agent must route critical escalations (I-9 Day 2+, background check flags) via Teams message or SMS — otherwise the 15-minute response SLA is structurally impossible by design, regardless of how good the agent logic is."),
    (ORANGE, "3", "Document every undocumented rule before launch",
     "Three known examples: (1) the SharePoint compliance matrix is outdated — Priya's desktop copy is the real version; (2) no Strategy department buddies for Finance hires — nowhere in writing; (3) Connections HR consistently mislabels hire_type — Priya manually corrects every one. If Priya leaves before these are captured, the knowledge leaves with her."),
]
for i, (color, num, title, body) in enumerate(steps):
    by = 1.42 + i * 1.78
    rect(s, 0.35, by, 0.62, 1.6, color)
    txt(s, num, 0.35, by+0.52, 0.62, 0.56, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.02, by, 12.0, 1.6, WHITE)
    txt(s, title, 1.18, by+0.1,  11.72, 0.42, size=13, bold=True, color=NAVY)
    txt(s, body,  1.18, by+0.52, 11.72, 1.0,  size=10.5, color=DARK)

rect(s, 0.35, 6.82, 12.6, 0.5, NAVY)
txt(s, "Phase 2 — when Agent 1 is proven:  Build Agent 2 (Proposal Router) for training track proposals, buddy matching, and IT access selection. Combined V×V score: 18 points — above the 15-point build threshold.",
    0.55, 6.86, 12.3, 0.42, size=10.5, color=WHITE)


# ── Save ─────────────────────────────────────────────────────────────────────
out = r"c:\Users\AnnOHagan\OneDrive - EPAM\FDE AI Program\FDE-practice-week2\FDE-practice-week2\ATX-Assessment-Summary.pptx"
prs.save(out)
print(f"Saved → {out}")

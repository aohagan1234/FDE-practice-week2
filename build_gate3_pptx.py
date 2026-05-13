"""Gate 3 presentation — MedFlex. Defense-focused, 10-minute coach call."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
MID    = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xF5, 0xF6, 0xF8)
ORANGE = RGBColor(0xE3, 0x6B, 0x2A)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
RED    = RGBColor(0xBF, 0x35, 0x2A)
AMBER  = RGBColor(0xCC, 0x84, 0x00)
LBLUE  = RGBColor(0xBB, 0xCC, 0xEE)
DBLUE  = RGBColor(0x23, 0x52, 0x96)
RULE   = RGBColor(0xDD, 0xDD, 0xDD)
LGREEN = RGBColor(0xE8, 0xF5, 0xEC)
LRED   = RGBColor(0xFB, 0xEB, 0xEA)
LAMBER = RGBColor(0xFD, 0xF5, 0xE0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 9


def slide():
    return prs.slides.add_slide(BLANK)

def box(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def t(s, text, x, y, w, h, size=15, bold=False, color=DARK,
      align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def header(s, title, subtitle=None):
    box(s, 0, 0, 13.33, 1.05, NAVY)
    t(s, title, 0.45, 0.08, 12.4, 0.62, size=26, bold=True, color=WHITE)
    if subtitle:
        t(s, subtitle, 0.45, 0.70, 12.4, 0.32, size=12, color=LBLUE, italic=True)

def slide_num(s, n):
    t(s, f"{n} / {TOTAL}", 12.5, 7.18, 0.75, 0.26, size=10, color=MID, align=PP_ALIGN.RIGHT)

def hline(s, y):
    box(s, 0.45, y, 12.43, 0.02, RULE)

def pill(s, text, x, y, w, h, bg, fg=WHITE):
    box(s, x, y, w, h, bg)
    t(s, text, x, y + 0.04, w, h - 0.04, size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, NAVY)
box(s, 0, 5.6, 13.33, 0.06, ORANGE)

t(s, "MEDFLEX  ·  GATE 3  ·  DEFENSE", 0.55, 0.38, 12.0, 0.45,
  size=12, bold=True, color=LBLUE)
t(s, "Candidate Ranking\n& Parallel Outreach",
  0.55, 0.9, 11.5, 2.1, size=42, bold=True, color=WHITE)
t(s, "Built, tested, and ready to defend.", 0.55, 3.1, 9.0, 0.45,
  size=16, color=LBLUE, italic=True)

# Three probe questions as callout boxes
probes = [
    (ORANGE, "Is this agentic — or just a\nmechanistic app with AI on top?"),
    (DBLUE,  "The CEO's two failed AI projects.\nHow is yours different?"),
    (RED,    "What kills this\nin production?"),
]
for i, (color, text) in enumerate(probes):
    bx = 0.55 + i * 4.28
    box(s, bx, 3.85, 4.0, 1.55, color)
    t(s, text, bx + 0.18, 3.98, 3.65, 1.28, size=14, bold=True, color=WHITE)

t(s, "Ann O'Hagan  ·  2026-05-13", 0.55, 5.78, 5.0, 0.38, size=11, color=LBLUE)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — WHAT WAS BUILT (SETUP)
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "What Was Built",
       "The brief, the two capabilities, and the wave structure.")
slide_num(s, 2)

# Left — the problem
t(s, "The problem", 0.45, 1.18, 5.8, 0.36, size=13, bold=True, color=NAVY)
hline(s, 1.54)

problem_lines = [
    "MedFlex places nurses into hospital shifts.",
    "When a shift needs filling, coordinators contact nurses one at a time.",
    "For urgent fills (≤4h), the window can expire before anyone confirms.",
    "Kim's team owns the hospital relationship — the agent never contacts\nthe hospital directly.",
]
for i, line in enumerate(problem_lines):
    by = 1.62 + i * 0.78
    box(s, 0.45, by + 0.18, 0.1, 0.1, NAVY)
    t(s, line, 0.68, by + 0.12, 5.55, 0.62, size=14, color=DARK)

box(s, 6.6, 1.12, 0.04, 5.9, RULE)

# Right — what was built
t(s, "What was built", 6.85, 1.18, 5.95, 0.36, size=13, bold=True, color=NAVY)
hline(s, 1.54)

capabilities = [
    (ORANGE, "Candidate ranking",
     "Ranks eligible nurses across 6 factors simultaneously.\n"
     "Returns up to 5 ranked candidates with a written rationale each."),
    (DBLUE,  "Parallel outreach",
     "Contacts top-ranked nurses simultaneously.\n"
     "Adapts strategy in real time based on urgency, pool size, and concurrent fills."),
]
for i, (color, title, body) in enumerate(capabilities):
    by = 1.62 + i * 1.75
    box(s, 6.85, by, 0.1, 1.45, color)
    t(s, title, 7.1, by + 0.1, 5.65, 0.42, size=14, bold=True, color=DARK)
    t(s, body,  7.1, by + 0.58, 5.65, 0.82, size=13, color=MID)

# Wave structure
by = 5.22
box(s, 0.45, by, 12.43, 0.36, NAVY)
t(s, "Wave structure", 0.65, by + 0.06, 12.0, 0.26, size=12, bold=True, color=WHITE)
box(s, 0.45, by + 0.36, 6.0, 1.55, LGREY)
t(s, "Wave 1 — 6-week pilot", 0.65, by + 0.44, 5.7, 0.36, size=13, bold=True, color=NAVY)
t(s, "Candidate ranking + parallel outreach.\nStandard fills >24h. 2–3 coordinators. Established hospitals.\nCoordinator approves every placement.",
  0.65, by + 0.82, 5.7, 0.98, size=12, color=DARK)
box(s, 6.45, by + 0.36, 6.43, 1.55, LGREY)
t(s, "Wave 2 — condition-triggered at month 5", 6.65, by + 0.44, 6.1, 0.36, size=13, bold=True, color=NAVY)
t(s, "Credential recency, no-show backfill, pre-shift mismatch detection.\nStarts only if first-recommendation acceptance rate ≥90% over 60 days.",
  6.65, by + 0.82, 6.1, 0.98, size=12, color=DARK)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — IS THIS AGENTIC?
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Is This Agentic — or Just a Fancy App?",
       "Honest answer: some parts are deterministic. Here is where the agent earns its place.")
slide_num(s, 3)

# Honest concession box
box(s, 0.45, 1.12, 12.43, 0.72, LGREY)
box(s, 0.45, 1.12, 0.12, 0.72, AMBER)
t(s, "The honest concession: urgency tier calculation, factor weights, no-show lookback, and response rate imputation "
     "are all deterministic — they are if/else rules and arithmetic. A script could run them.",
  0.72, 1.2, 12.0, 0.56, size=13, color=DARK, italic=True)

# Three columns: deterministic / where agent adds value / the test
cols = [
    (RULE, DARK, "What a standard app does",
     [
         "Contact nurses sequentially —\none at a time",
         "Apply a fixed contact order\n(e.g. alphabetical or seniority)",
         "Stop when someone confirms\nor all are exhausted",
         "No explanation of why\na nurse was ranked where they were",
         "Fail silently when data\nis missing or a system is down",
     ]),
    (NAVY, WHITE, "What the agent does differently",
     [
         "Contacts multiple nurses simultaneously,\nadapting N to urgency and pool size",
         "Ranks by 6 factors with different\nweights depending on urgency context",
         "Coordinates across concurrent fills —\nexcludes nurses already in active outreach",
         "Generates a specific written rationale\nper candidate that the coordinator reads",
         "Degrades gracefully with imputed data\nand flags what changed in the rationale",
     ]),
    (GREEN, WHITE, "The test: could a fixed rule do it?",
     [
         "No — N changes based on live\npool state + urgency at ranking time",
         "No — weight combination varies\nby urgency tier AND specialist flag",
         "No — requires real-time shared\nstate across concurrent fill cycles",
         "No — rationale is specific to\nthis candidate in this context",
         "No — a script errors or silently\nwrong; agent explains and escalates",
     ]),
]

for i, (hdr_color, hdr_text_color, title, items) in enumerate(cols):
    cx = 0.45 + i * 4.3
    cw = 4.1
    box(s, cx, 1.95, cw, 0.38, hdr_color if hdr_color != RULE else DBLUE)
    t(s, title, cx + 0.1, 1.99, cw - 0.15, 0.3,
      size=11, bold=True, color=hdr_text_color if hdr_color != RULE else WHITE)
    for j, item in enumerate(items):
        by = 2.38 + j * 0.94
        box(s, cx, by, cw, 0.9, LGREY if j % 2 == 0 else WHITE)
        t(s, item, cx + 0.15, by + 0.08, cw - 0.25, 0.76, size=11, color=DARK)

box(s, 0.45, 7.08, 12.43, 0.36, GREEN)
t(s, "The agent value is the COMBINATION: simultaneous outreach + contextual ranking + real-time state awareness + explainable output. "
     "No single rule handles all four.",
  0.65, 7.13, 12.0, 0.26, size=11, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — HOW IS THIS DIFFERENT FROM FAILED AI PROJECTS?
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "The CEO's Two Failed AI Projects — How Is This Different?",
       "AI projects fail in predictable ways. Here is how each one was addressed.")
slide_num(s, 4)

t(s, "Common failure", 0.45, 1.12, 5.8, 0.32, size=12, bold=True, color=RED)
t(s, "What MedFlex does instead", 7.1, 1.12, 5.95, 0.32, size=12, bold=True, color=GREEN)
hline(s, 1.44)

failures = [
    ("Automated the wrong thing — solved a problem nobody had",
     "The bottleneck is real and measured: coordinators spend fill windows on sequential manual contact. "
     "The two capabilities target fill time directly."),

    ("Black box — nobody could see why it made decisions",
     "Every ranked candidate has a written rationale the coordinator reads before approving. "
     "The coordinator sees why, can override, and that override is logged."),

    ("Removed humans from decisions they needed to own",
     "The coordinator approves every placement. The agent never contacts the hospital. "
     "Kim's team owns the hospital relationship — the agent does not touch it."),

    ("Launched too big — no recovery when it went wrong",
     "Wave 1 is 2–3 coordinators, standard fills only, established accounts. "
     "It is small enough to stop and reverse without damaging live operations."),

    ("No success criteria — nobody knew if it was working",
     "First-recommendation acceptance rate ≥90% over 60 days triggers Wave 2. "
     "If that threshold is not hit, Wave 2 does not start. The test is defined before build."),
]

for i, (fail, fix) in enumerate(failures):
    by = 1.52 + i * 1.14
    bg = LGREY if i % 2 == 0 else WHITE
    box(s, 0.45, by, 12.43, 1.1, bg)
    box(s, 0.45, by, 0.12, 1.1, RED)
    box(s, 6.82, by, 0.04, 1.1, RULE)
    box(s, 6.86, by, 0.12, 1.1, GREEN)
    t(s, fail, 0.72, by + 0.12, 5.9, 0.86, size=12, bold=True, color=DARK)
    t(s, fix,  7.14, by + 0.12, 5.9, 0.86, size=12, color=MID)


# ═══════════════════════════════════════════════════════════
# SLIDE 5 — WHAT KILLS THIS IN PRODUCTION?
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "What Kills This in Production?",
       "Three honest failure modes — identified before build, not after.")
slide_num(s, 5)

t(s, "These were identified in the validation plan. Two have designed mitigations. One is the real risk.",
  0.45, 1.14, 12.43, 0.36, size=13, color=DARK, italic=True)

risks = [
    (AMBER, "MEDIUM", "Concurrent fill state inconsistency",
     "If two fills run simultaneously and shared state has a stale read, a nurse already committed to one fill "
     "gets contacted for a second. The outreach window for the second fill is wasted.",
     "Designed for: fallback to database flag per nurse if shared state unavailable. "
     "Coordinator is notified when state could not be confirmed. Not fully testable before live volume."),

    (AMBER, "MEDIUM", "Nurse confirms via a different channel than contacted",
     "Nurse is contacted by SMS and calls the coordinator directly. The system does not detect the confirmation. "
     "Outreach window expires. Next candidate is contacted. If the first nurse shows up, the shift may be double-placed.",
     "Designed for: noted as a known production failure mode in the validation plan. "
     "Cannot be tested before live operation — depends on nurse behaviour patterns."),

    (RED, "HIGH", "CRM data sparsity — the silent killer",
     "If coordinators have not consistently logged response rates and no-show history, "
     "the ranking agent operates on credentials and proximity only. "
     "The rationale thins. Coordinators increasingly override. Acceptance rate never reaches 90%. Wave 2 never starts.",
     "This is the most likely real failure. The system works correctly by its own logic "
     "while producing rankings coordinators do not trust. It fails silently."),
]

for i, (color, level, title, problem, mitigation) in enumerate(risks):
    by = 1.60 + i * 1.85
    box(s, 0.45, by, 12.43, 1.78, LGREY if i % 2 == 0 else WHITE)
    box(s, 0.45, by, 0.12, 1.78, color)
    pill(s, level, 0.7, by + 0.12, 1.1, 0.38, color)
    t(s, title, 1.95, by + 0.12, 10.75, 0.38, size=14, bold=True, color=DARK)
    t(s, problem, 0.72, by + 0.58, 6.1, 1.1, size=12, color=DARK)
    box(s, 6.88, by + 0.5, 0.04, 1.2, RULE)
    t(s, mitigation, 7.08, by + 0.58, 5.65, 1.1, size=11, color=MID, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 6 — MARCUS'S FEEDBACK AND WHAT CHANGED
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Client Feedback — What Marcus Raised and What Changed",
       "Three pushbacks. Each one addressed. One was already in the original design.")
slide_num(s, 6)

t(s, "Feedback", 0.45, 1.12, 5.8, 0.32, size=12, bold=True, color=ORANGE)
t(s, "What changed", 7.1, 1.12, 5.95, 0.32, size=12, bold=True, color=GREEN)
hline(s, 1.44)

feedback_rows = [
    ("Board meeting is in 6 weeks, not 8 — restructure the plan",
     "Plan restructured to 6 weeks. Wave 1 — parallel outreach + coordinator approval gate — running on 20–30 live shifts "
     "by week 6. Board deliverable is fill time and first-recommendation acceptance rate data. "
     "Not proof of 10x capacity — that is the honest version of what 6 weeks can demonstrate."),

    ("Remove credential verification from Wave 1",
     "Credential verification moved to Wave 2. Compliance team continues quarterly cadence during the pilot. "
     "Worth noting: the original design already flagged credential verification as conditional on Linda confirming "
     "the root cause of the 7% mismatch rate. Removing it from Wave 1 is consistent with that — not a concession."),

    ("Kim flagged: who manages the hospital when things go wrong? The spec did not say.",
     "Failure path section added to the architecture. Four scenarios covered: no confirmation within the outreach window, "
     "pool exhausted, parsing error at coordinator approval gate, no-show detected post-placement. "
     "Coordinator owns the hospital relationship in all four cases. Agent pre-drafts the status message "
     "and surfaces the remaining pool state — coordinator makes the call within 15–30 minutes of escalation."),
]

for i, (feedback, change) in enumerate(feedback_rows):
    by = 1.52 + i * 1.9
    bg = LGREY if i % 2 == 0 else WHITE
    box(s, 0.45, by, 12.43, 1.82, bg)
    box(s, 0.45, by, 0.12, 1.82, ORANGE)
    box(s, 6.82, by, 0.04, 1.82, RULE)
    box(s, 6.86, by, 0.12, 1.82, GREEN)
    t(s, feedback, 0.72, by + 0.16, 5.9, 1.5, size=12, bold=True, color=DARK)
    t(s, change,   7.14, by + 0.16, 5.9, 1.5, size=12, color=MID)


# ═══════════════════════════════════════════════════════════
# SLIDE 7 — BUILD LOOP LEARNINGS
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Build Loop — What Building Against My Own Spec Revealed",
       "Five signals. Four were gaps in my spec. One was a decision the builder made that I never asked for.")
slide_num(s, 7)

# Key observation
box(s, 0.45, 1.12, 12.43, 0.72, LGREY)
box(s, 0.45, 1.12, 0.12, 0.72, AMBER)
t(s, "No clarifying questions were asked during the build. Wherever my spec was unclear, the builder resolved it silently. "
     "I only found out what decisions had been made by reading the output. "
     "A builder who asks nothing either has a perfect spec or has made undocumented assumptions.",
  0.72, 1.18, 12.0, 0.58, size=13, color=DARK, italic=True)

# Table
box(s, 0.45, 1.94, 12.43, 0.38, NAVY)
t(s, "What happened",             0.65, 2.00, 5.9,  0.26, size=11, bold=True, color=WHITE)
t(s, "Classification",            6.72, 2.00, 2.65, 0.26, size=11, bold=True, color=WHITE)
t(s, "What my spec failed to say",9.52, 2.00, 3.2,  0.26, size=11, bold=True, color=WHITE)

signal_rows = [
    (RED,   "Builder invented a proximity formula using a reference distance I never specified",
             "Unjustified\nimplementation choice",
             "I listed proximity as a factor. I did not say how to turn distance into a number."),
    (AMBER, '"Within 10%" applied to the overall composite score, not to each factor individually',
             "Spec gap",
             "My wording sounded specific. It did not say what was being compared."),
    (AMBER, "CRM unavailable treated the same as data simply missing for a nurse",
             "Spec gap",
             "No mechanism to distinguish system down from data not existing for that nurse."),
    (AMBER, "System unavailability rules could not be implemented as written",
             "Spec gap",
             "The tool had no way of being told a system was unavailable. The rule had no mechanism."),
    (AMBER, "pool_exhausted escalation had no field for the excluded candidates list",
             "Spec gap",
             "I said include it. The output type I designed had nowhere to put it."),
]

row_bgs = [WHITE, LGREY, WHITE, LGREY, WHITE]
for i, (color, signal, classification, failure) in enumerate(signal_rows):
    ry = 2.32 + i * 0.94
    box(s, 0.45, ry, 12.43, 0.9, row_bgs[i])
    box(s, 0.45, ry, 0.12,  0.9, color)
    t(s, signal,         0.68, ry + 0.12, 5.85, 0.66, size=12, color=DARK)
    t(s, classification, 6.72, ry + 0.12, 2.6,  0.66, size=12, color=DARK, bold=True)
    t(s, failure,        9.52, ry + 0.12, 3.15, 0.66, size=11, color=MID, italic=True)

box(s, 0.45, 7.04, 12.43, 0.38, GREEN)
t(s, "All five addressed in the updated D4a spec: scoring formula, rule 7 threshold defined, "
     "system availability inputs added, excluded_candidates field defined.",
  0.65, 7.08, 12.0, 0.28, size=11, bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 8 — REFLECTION AND LEARNINGS
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "Reflection — What I Would Do Differently",
       "Two from how I ran the engagement. Four from how I wrote the spec.")
slide_num(s, 8)

box(s, 0.45, 1.12, 5.95, 0.36, NAVY)
t(s, "Engagement", 0.65, 1.16, 5.65, 0.28, size=12, bold=True, color=WHITE)

eng = [
    ("Confirm the timeline in the first conversation",
     "Marcus told me the board meeting was in 6 weeks, not 8, after I had already structured the plan. "
     "That is a basic question I should have asked before writing anything."),
    ("Involve operations (Kim) in discovery — not just leadership",
     "Kim was not in the discovery session. The failure path gap she spotted would have been caught "
     "in a 30-minute conversation I simply did not have."),
]
for i, (title, body) in enumerate(eng):
    by = 1.52 + i * 1.55
    box(s, 0.45, by, 0.12, 1.38, NAVY)
    t(s, title, 0.72, by + 0.08, 5.55, 0.38, size=13, bold=True, color=DARK)
    t(s, body,  0.72, by + 0.52, 5.55, 0.78, size=12, color=MID)

box(s, 6.6, 1.08, 0.04, 5.9, RULE)

box(s, 6.86, 1.12, 5.92, 0.36, ORANGE)
t(s, "Spec writing", 7.06, 1.16, 5.62, 0.28, size=12, bold=True, color=WHITE)

spec = [
    ("Define the output before writing the error handling rule",
     "If the output has no field for something the rule requires, the rule cannot be delivered. "
     "I wrote 'include excluded candidates in the escalation' — the output had nowhere to put them."),
    ("For every 'if X is unavailable' rule — ask how the tool finds out",
     "I wrote rules for system unavailability without thinking about what signal the tool would receive. "
     "A rule that depends on information never passed in is not a rule."),
    ("Every ranking factor needs a formula, not just a name",
     "I listed proximity as a factor. The builder invented the formula. "
     "Different formulas produce different rankings from the same pool."),
    ("'Within 10%' sounds specific — it is not",
     "I needed to say: compare the overall scores, and boost only if the gap is 10% or less. "
     "The builder had to guess what was being compared."),
]
for i, (title, body) in enumerate(spec):
    by = 1.52 + i * 1.45
    box(s, 6.86, by, 0.12, 1.28, ORANGE)
    t(s, title, 7.14, by + 0.08, 5.5, 0.38, size=13, bold=True, color=DARK)
    t(s, body,  7.14, by + 0.52, 5.5, 0.7,  size=12, color=MID)


# ═══════════════════════════════════════════════════════════
# SLIDE 9 — SUMMARY: THE HONEST TRADE-OFFS
# ═══════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, 13.33, 7.5, WHITE)
header(s, "The Honest Trade-offs",
       "What we are confident about. What we are not. What would change the design.")
slide_num(s, 9)

# Three columns
cols_data = [
    (GREEN, LGREEN, "Confident",
     [
         "Fill time reduction is measurable — first-recommendation acceptance rate is a real metric",
         "Coordinator trust is built incrementally — rationale + approval gate + override logging",
         "Agent never touches the hospital relationship",
         "Wave 2 is gated — does not start unless the pilot data supports it",
         "Production failure modes are identified and designed for before build",
     ]),
    (AMBER, LAMBER, "Uncertain",
     [
         "CRM data quality — if response rates and no-show history are sparse, ranking quality is low",
         "Urgency tier thresholds — the 4h/24h boundaries are configurable defaults, not validated against MedFlex data",
         "in_active_outreach state consistency at concurrent fill volume not testable before live operation",
         "Nurse confirmation via different channel — a known gap with no technical fix",
     ]),
    (RED, LRED, "What changes the design",
     [
         "If CRM data is too sparse: Wave 1 starts with credential + proximity ranking only — rationale flags this explicitly",
         "If shared state infrastructure doesn't exist: database flag per nurse (slower, same outcome)",
         "If urgency thresholds don't match MedFlex patterns: reconfigure after weeks 2–5 pilot data",
         "If acceptance rate stays below 90%: investigate — ranking or outreach, not both at once",
     ]),
]

for i, (hdr_color, bg_color, title, items) in enumerate(cols_data):
    cx = 0.45 + i * 4.3
    cw = 4.08
    box(s, cx, 1.12, cw, 0.42, hdr_color)
    t(s, title, cx + 0.15, 1.16, cw - 0.2, 0.34,
      size=14, bold=True, color=WHITE)
    box(s, cx, 1.54, cw, 5.6, bg_color)
    for j, item in enumerate(items):
        box(s, cx + 0.15, 1.68 + j * 1.32, 0.1, 0.1, hdr_color)
        t(s, item, cx + 0.38, 1.62 + j * 1.32, cw - 0.48, 1.2,
          size=11, color=DARK)


# ── Save ──────────────────────────────────────────────────
OUT = (r"C:\Users\AnnOHagan\OneDrive - EPAM\FDE AI Program"
       r"\FDE-practice-week2\FDE-practice-week2"
       r"\Week3_Docs\Gate3-MedFlex-Defense-v2.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")

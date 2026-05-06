"""Generate Gate 2 presentation — Apex Distribution Ltd, ETA Inquiry Agent."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (matches build_pptx.py) ──────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
ORANGE = RGBColor(0xE3, 0x6B, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
MID    = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xF2, 0xF4, 0xF7)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
RED    = RGBColor(0xBF, 0x35, 0x2A)
AMBER  = RGBColor(0xCC, 0x84, 0x00)
TEAL   = RGBColor(0x0E, 0x7C, 0x7B)
LBLUE  = RGBColor(0xBB, 0xCC, 0xEE)
DBLUE2 = RGBColor(0x23, 0x52, 0x96)
LGREEN = RGBColor(0xD6, 0xEE, 0xDC)
LRED   = RGBColor(0xF7, 0xDC, 0xDA)
LAMBER = RGBColor(0xFD, 0xF0, 0xD5)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def mtxt(s, x, y, w, h, lines):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align, italic) in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.italic = italic
    return tb

def header(s, title, subtitle=None):
    rect(s, 0, 0, 13.33, 1.25, NAVY)
    txt(s, title, 0.38, 0.07, 12.5, 0.7, size=26, bold=True, color=WHITE)
    if subtitle:
        txt(s, subtitle, 0.38, 0.77, 12.5, 0.42, size=12, color=LBLUE)

def slide_number(s, n, total=6):
    txt(s, f"{n} / {total}", 12.6, 7.15, 0.65, 0.3,
        size=9, color=MID, align=PP_ALIGN.RIGHT)

def divider(s, y):
    rect(s, 0.35, y, 12.63, 0.03, RGBColor(0xCC, 0xCC, 0xCC))


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 4.8,  13.33, 2.7,  DBLUE2)
rect(s, 0, 2.82, 13.33, 0.07, ORANGE)

txt(s, "APEX DISTRIBUTION LTD  ·  GATE 2  ·  CUSTOMER OPERATIONS",
    0.55, 0.32, 12.2, 0.42, size=10, color=LBLUE, bold=True)
txt(s, "ETA Inquiry Agent\nBefore. After. Why.",
    0.55, 0.82, 12.0, 1.9, size=42, bold=True, color=WHITE)
txt(s, "A 4-minute analysis of W2 — what 400 inquiries a day actually costs, "
       "what the redesigned process looks like, and why each delegation call is safe.",
    0.55, 2.95, 11.5, 0.6, size=13, color=LBLUE)

# Left callout — the number
rect(s, 0.55, 5.0, 4.5, 1.72, ORANGE)
txt(s, "26.7 person-hours / day",
    0.72, 5.08, 4.2, 0.52, size=18, bold=True, color=WHITE)
txt(s, "400 inquiries × 4 min avg\n= 3.3 FTE doing database lookups",
    0.72, 5.56, 4.2, 1.05, size=13, color=WHITE)

# Right callout — the target
mtxt(s, 5.35, 5.12, 7.3, 1.45, [
    ("Target:  agent handles >80% autonomously", 13, True,  LBLUE, PP_ALIGN.LEFT, False),
    ("2.8 FTE freed  ·  1–2 month payback  ·  6 weeks to production",
     12, False, LBLUE, PP_ALIGN.LEFT, False),
    ("",  6, False, WHITE, PP_ALIGN.LEFT, False),
    ("Presenter: Ann O'Hagan  ·  2026-05-06",
     10, False, RGBColor(0x88,0x99,0xBB), PP_ALIGN.LEFT, False),
])


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — BEFORE: THE CURRENT PROCESS
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "Before — The Current Process",
       "Where the time goes, where the cognitive load sits, where it breaks")
slide_number(s, 2)

# Zone A box
rect(s, 0.35, 1.38, 5.85, 0.42, NAVY)
txt(s, "Zone A — Status lookup  (~1 min, deterministic)",
    0.5, 1.41, 5.6, 0.36, size=12, bold=True, color=WHITE)
rect(s, 0.35, 1.80, 5.85, 2.02, WHITE)
mtxt(s, 0.5, 1.85, 5.55, 1.9, [
    ("Agent looks up consignment in Salesforce CRM.", 11, False, DARK, PP_ALIGN.LEFT, False),
    ("Reads a status: pre-dispatch / out-for-delivery / delivered / exception.", 11, False, DARK, PP_ALIGN.LEFT, False),
    ("", 5, False, DARK, PP_ALIGN.LEFT, False),
    ("No judgment required. This is a database query.", 11, True, NAVY, PP_ALIGN.LEFT, False),
    ("", 5, False, DARK, PP_ALIGN.LEFT, False),
    ("Data source: Salesforce CRM REST API  ·  5–10 min sync lag",
     10, False, MID, PP_ALIGN.LEFT, True),
])

# Zone B box
rect(s, 0.35, 3.96, 5.85, 0.42, RED)
txt(s, "Zone B — ETA decision & driver sync  (1–3 min, where load sits)",
    0.5, 3.99, 5.6, 0.36, size=12, bold=True, color=WHITE)
rect(s, 0.35, 4.38, 5.85, 2.87, WHITE)
mtxt(s, 0.5, 4.43, 5.55, 2.75, [
    ("If status = out-for-delivery, agent needs a tighter ETA.", 11, False, DARK, PP_ALIGN.LEFT, False),
    ("Today: manually contact dispatch → dispatch contacts driver → wait.", 11, False, DARK, PP_ALIGN.LEFT, False),
    ("", 5, False, DARK, PP_ALIGN.LEFT, False),
    ("Artefact 3 evidence:", 11, True, NAVY, PP_ALIGN.LEFT, False),
    ("Customer asked at 11:14. Agent checked dispatch at 11:19.", 11, False, DARK, PP_ALIGN.LEFT, False),
    ("Best-guess reply sent at 11:24. 10 minutes. Manual GPS ping.", 11, False, DARK, PP_ALIGN.LEFT, False),
    ("Customer got a 4-hour window and wanted something tighter.", 11, False, DARK, PP_ALIGN.LEFT, False),
    ("", 5, False, DARK, PP_ALIGN.LEFT, False),
    ("Driver non-response is routine — not an edge case.", 11, True, RED, PP_ALIGN.LEFT, False),
])

# Right: 4 failure modes
rect(s, 6.45, 1.38, 6.53, 0.42, ORANGE)
txt(s, "Where it breaks — 4 failure modes",
    6.6, 1.41, 6.28, 0.36, size=12, bold=True, color=WHITE)

failures = [
    (RED,   "CRM lag",
     "Status is 5–10 min behind reality. Driver may have already delivered when agent says 'out-for-delivery.'"),
    (RED,   "No deduplication",
     "Same customer asks twice in an hour → two separate driver queries. No memory across inquiries."),
    (AMBER, "4-hour window dissatisfies",
     "The SOP answer (13:00–17:00) frustrates customers. Tighter ETA requires the 10-minute manual process above."),
    (AMBER, "SOP is stale",
     "Section 4.2 references DispatchHub — retired October 2024. Section 4.3 says 'TBD pending review.' Not updated since Oct 2023."),
]
for i, (color, title, body) in enumerate(failures):
    by = 1.88 + i * 1.3
    rect(s, 6.45, by, 0.38, 1.18, color)
    rect(s, 6.83, by, 6.15, 1.18, WHITE)
    txt(s, title, 6.98, by + 0.08, 5.88, 0.36, size=11, bold=True, color=DARK)
    txt(s, body,  6.98, by + 0.42, 5.88, 0.68, size=10, color=MID)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AFTER: THE REDESIGNED PROCESS
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "After — The Redesigned Process",
       "Three layers. Explicit about which steps are agentic, which are hybrid, which are human.")
slide_number(s, 3)

cols = [
    (GREEN, LGREEN, "WORKFLOW LAYER\nFully agentic", [
        "Consignment ID validation",
        "CRM lookup (Salesforce REST API)",
        "Status = Delivered → template reply",
        "Status = Pre-dispatch → template reply",
        "Status = Exception → escalate to dispatcher",
        "Status = Null → escalate (data error)",
        "Inquiry deduplication (1h cache)",
        "Detect explicit tight-ETA request\n(keyword match — deterministic rule)",
    ], "~80% of all inquiries close here\nNo LLM. No judgment. A rule engine."),

    (AMBER, LAMBER, "AGENT LAYER\nAgent-led, human on exceptions", [
        "Contact driver via driver app",
        "Wait up to 2 min for response",
        "GPS mode → calculate tight ETA window",
        "MSG mode → confirm 'on track' + fallback",
        "No response → send fallback 13:00–17:00",
        "Explicit tight-ETA request → escalate\nto dispatcher (10% optional path)",
    ], "Driver contact is async + best-effort.\nFallback is primary path, not failure mode."),

    (RED, LRED, "HUMAN LAYER\nHuman-led, agent surfaces data", [
        "Consignment not found in CRM\n→ manual records check (paper/email)",
        "Exception flag\n→ dispatcher decides next step",
        "Null status\n→ IT investigates CRM sync issue",
    ], "Not because the lookup is hard.\nBecause what happens NEXT\nrequires dispatcher judgment."),
]

for i, (hdr_color, body_color, title, items, footer) in enumerate(cols):
    cx = 0.35 + i * 4.34
    cw = 4.15
    rect(s, cx, 1.38, cw, 0.72, hdr_color)
    txt(s, title, cx + 0.12, 1.41, cw - 0.2, 0.66,
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, cx, 2.10, cw, 4.3, body_color)
    for j, item in enumerate(items):
        txt(s, f"• {item}", cx + 0.18, 2.18 + j * 0.5, cw - 0.3, 0.46,
            size=10, color=DARK)
    rect(s, cx, 6.40, cw, 0.82, hdr_color)
    txt(s, footer, cx + 0.12, 6.44, cw - 0.2, 0.74,
        size=9.5, color=WHITE, italic=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHY: DEFENDING THE DELEGATION CALLS
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "Why — Defending the 3 Delegation Calls",
       "Each boundary is a deliberate choice, not a default. Here is the rationale.")
slide_number(s, 4)

calls = [
    (GREEN, "Driver non-response → fallback is FULLY AGENTIC",
     "Claim:", "The fallback is a timeout rule, not a runtime decision.",
     "Because if two minutes pass with no driver response, send the template. "
     "The trigger is a clock. The message is pre-approved. There is no judgment — "
     "the human made the judgment when they approved the protocol. "
     "Moving this to human-in-the-loop adds 3–5 minutes and breaks the SLA.",
     "Risk if wrong:", "Stale ETA sent. Customer gets 4h window when driver was 10 min away. "
     "Reputational, not regulatory. Accepted as a known, bounded trade-off."),

    (AMBER, "Consignment lookup → WORKFLOW ENGINE, not LLM",
     "Claim:", "This is a switch statement on 5 enum values. LLM adds zero value here.",
     "CRM returns a status: pre-dispatch / out-for-delivery / delivered / exception / null. "
     "The agent selects a template. No natural language reasoning. No ambiguity. "
     "Hallucination risk on an LLM call is real; it is zero on a database query. "
     "An event-driven API wrapper is faster, cheaper, and strictly safer.",
     "Risk if wrong:", "None — if the CRM returns a status the workflow doesn't recognise, "
     "it escalates to human rather than guessing. The failure mode is safe by design."),

    (RED, "Exception flag → HUMAN-LED",
     "Claim:", "Not because the lookup is hard — because the next step requires judgment.",
     "The exception reason is in CRM. The agent can read it. What the agent cannot do is "
     "decide what happens next: retry tomorrow? Insurance claim? Goodwill credit? "
     "Those decisions require knowledge of customer history, dispatcher availability, and "
     "insurance pre-auth that no deterministic rule can encode. "
     "Agent surfaces the data. Dispatcher makes the call.",
     "Risk if wrong:", "High. Exception handling is the work stream most likely to escalate "
     "into a customer complaint or contract dispute. Keeping the human in the loop here "
     "is not a failure of ambition — it is the correct delegation call."),
]

for i, (color, title, lbl1, claim, body, lbl2, risk) in enumerate(calls):
    by = 1.38 + i * 1.97
    rect(s, 0.35, by, 0.22, 1.83, color)
    rect(s, 0.57, by, 12.41, 0.42, color)
    txt(s, title, 0.72, by + 0.05, 12.12, 0.34, size=12, bold=True, color=WHITE)
    rect(s, 0.57, by + 0.42, 12.41, 1.41, WHITE)
    txt(s, lbl1,  0.72, by + 0.48, 1.0,   0.30, size=10, bold=True,  color=DARK)
    txt(s, claim, 1.72, by + 0.48, 5.5,   0.30, size=10, bold=True,  color=NAVY)
    txt(s, body,  0.72, by + 0.78, 7.5,   0.95, size=10, color=DARK)
    txt(s, lbl2,  8.42, by + 0.48, 1.35,  0.30, size=10, bold=True,  color=DARK)
    txt(s, risk,  8.42, by + 0.78, 4.4,   0.95, size=10, color=MID, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — KEY ASSUMPTIONS & OPEN QUESTIONS
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "Assumptions & Open Questions",
       "What the design depends on — and what would change if the answer is different")
slide_number(s, 5)

txt(s, "Three questions must be answered before Phase 1 build starts. "
       "Each one changes the architecture, not just a configuration value.",
    0.38, 1.33, 12.55, 0.48, size=12, color=DARK, italic=True)

questions = [
    (RED, "Q1 — Driver app API surface",
     "LOW confidence",
     "GPS API exists",
     "Agent calculates tight ETA in <30 sec. 2-min wait protocol is unnecessary. Response time improves.",
     "Message-only",
     "2-min wait + fallback is correct. Keep spec as is. SLA floor is ~3 min, not ~1 min.",
     "No API at all",
     "Agent becomes a CRM lookup tool only. Significant capability reduction. Phase 1 scope changes."),

    (AMBER, "Q2 — CRM status freshness",
     "MEDIUM confidence",
     "Real-time (<1 min)",
     "Agent can trust CRM status. ETA accuracy limited only by driver location freshness.",
     "5–30 min polling",
     "Agent must include staleness disclaimer in every response. ±15 min accuracy target may need widening.",
     "End-of-shift batch",
     "Core ETA logic breaks — agent cannot distinguish 'in transit' from 'delivered but not updated.' Full redesign required."),

    (ORANGE, "Q3 — Sandra credit audit gap",
     "LOW confidence",
     "Known gap — overrides bypass export",
     "APEX_CREDITS is incomplete. Phase 2 billing agent will misidentify open disputes. Compliance exposure.",
     "Off-policy (Sandra shouldn't)",
     "Audit trail is theoretically intact. Phase 2 agent enforces structured path — closes gap automatically.",
     "Export is complete",
     "Internal note is wrong or refers to a lag. Lower risk — timing issue only."),
]

for i, (color, title, conf, a_hdr, a_body, b_hdr, b_body, c_hdr, c_body) in enumerate(questions):
    by = 1.90 + i * 1.72
    rect(s, 0.35, by, 3.72, 0.38, color)
    txt(s, title, 0.50, by + 0.04, 3.4, 0.30, size=11, bold=True, color=WHITE)
    txt(s, conf,  4.12, by + 0.04, 1.8, 0.30, size=10, color=color, bold=True)

    for j, (hdr, body) in enumerate([(a_hdr, a_body), (b_hdr, b_body), (c_hdr, c_body)]):
        bx = 0.35 + j * 4.34
        rect(s, bx, by + 0.38, 4.15, 0.28, RGBColor(0xDD, 0xDD, 0xDD))
        txt(s, f"If: {hdr}", bx + 0.12, by + 0.40, 3.9, 0.24, size=9.5, bold=True, color=DARK)
        rect(s, bx, by + 0.66, 4.15, 0.96, WHITE)
        txt(s, body, bx + 0.12, by + 0.70, 3.9, 0.88, size=9.5, color=MID, italic=True)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SUMMARY & WHAT'S NEEDED
# ═════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
header(s, "Summary", "The recommendation in three lines — and what we need from you")
slide_number(s, 6)

# Before/After summary table
rect(s, 0.35, 1.38, 12.63, 0.40, NAVY)
txt(s, "Before",  0.50,  1.42, 4.0, 0.32, size=12, bold=True, color=WHITE)
txt(s, "After",   5.00,  1.42, 4.0, 0.32, size=12, bold=True, color=WHITE)
txt(s, "Change",  9.55,  1.42, 3.3, 0.32, size=12, bold=True, color=WHITE)

rows = [
    ("Volume handled",         "400/day, all manual",     "320+/day autonomous (>80%)",    "−80% manual load"),
    ("Avg response time",      "4 min (P50)",             "<3 min (P50) target",            "Faster on most paths"),
    ("Driver contact",         "Manual via dispatch",     "Direct via driver app (async)",  "Removes dispatch as bottleneck"),
    ("Duplicate queries",      "No deduplication",        "1h cache — same inquiry = cache","Eliminates duplicate driver pings"),
    ("FTE impact",             "3.3 FTE on ETA alone",    "2.8 FTE freed for redeployment", "£140K/month redeployment value"),
    ("ROI payback",            "N/A (no agent)",          "1–2 months",                     "Low build risk"),
]
for i, (label, before, after, change) in enumerate(rows):
    by = 1.78 + i * 0.54
    bg = WHITE if i % 2 == 0 else LGREY
    rect(s, 0.35, by, 12.63, 0.54, bg)
    txt(s, label,  0.50,  by + 0.08, 4.0, 0.36, size=10.5, bold=True,  color=DARK)
    txt(s, before, 0.50,  by + 0.08, 4.3, 0.36, size=10.5, bold=False, color=MID)
    txt(s, after,  5.00,  by + 0.08, 4.3, 0.36, size=10.5, bold=False, color=GREEN)
    txt(s, change, 9.55,  by + 0.08, 3.3, 0.36, size=10.5, bold=True,  color=NAVY)

# What we need
rect(s, 0.35, 5.12, 12.63, 0.40, ORANGE)
txt(s, "Three things needed before build starts",
    0.50, 5.15, 12.3, 0.32, size=12, bold=True, color=WHITE)

needs = [
    ("Q1", "Driver app API surface — GPS, message-only, or no API? Controls the entire Phase 1 architecture."),
    ("Q2", "CRM status freshness — seconds, minutes, or batch? Changes the accuracy claim and staleness disclaimer."),
    ("Q3", "Credits audit trail gap — systemic or one-off? Determines whether APEX_CREDITS is a reliable data source for Phase 2."),
]
for i, (label, body) in enumerate(needs):
    bx = 0.35 + i * 4.22
    rect(s, bx, 5.52, 0.52, 1.75, ORANGE)
    txt(s, label, bx, 5.52, 0.52, 1.75,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, bx + 0.52, 5.52, 3.68, 1.75, WHITE)
    txt(s, body, bx + 0.62, 5.60, 3.48, 1.6, size=10, color=DARK)


# ── Save ─────────────────────────────────────────────────────────────────────
OUT = (r"C:\Users\AnnOHagan\OneDrive - EPAM\FDE AI Program"
       r"\FDE-practice-week2\FDE-practice-week2"
       r"\FDE-assessment-week2\Gate2-Presentation.pptx")
prs.save(OUT)
print(f"Saved: {OUT}")

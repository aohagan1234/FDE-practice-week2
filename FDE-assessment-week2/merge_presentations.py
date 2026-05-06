"""Merge and streamline Apex presentations - keep only key points"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
ORANGE = RGBColor(0xE3, 0x6B, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
MID    = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xF2, 0xF4, 0xF7)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
RED    = RGBColor(0xBF, 0x35, 0x2A)
AMBER  = RGBColor(0xCC, 0x84, 0x00)
LAMBER = RGBColor(0xFD, 0xF0, 0xD5)
LRED   = RGBColor(0xF7, 0xDC, 0xDA)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet_list(s, items, x, y, w, h, size=16, color=DARK):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.level = 0
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
rect(s, 0, 0, 13.33, 1.5, NAVY)
txt(s, "Apex Distribution Ltd", 0.5, 0.4, 12, 0.7, size=40, bold=True, color=WHITE)
txt(s, "Gate 2 Assessment: Build ETA Inquiry Agent First", 0.5, 3.2, 12, 0.5, size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM + DECISION
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "Problem & Solution", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Problem
txt(s, "The Problem", 0.8, 1.5, 5, 0.3, size=18, bold=True, color=NAVY)
bullet_list(s, [
    "35 staff handling 730 requests/day",
    "4 work streams analyzed",
    "2 prior automation failures"
], 0.8, 1.9, 5, 1.5, size=15)

# Solution
txt(s, "The Solution", 7.3, 1.5, 5, 0.3, size=18, bold=True, color=GREEN)
bullet_list(s, [
    "Build W2: ETA Agent first",
    "400 inquiries/day automated",
    "1-2 month ROI payback",
    "Low risk, high volume"
], 7.3, 1.9, 5, 1.5, size=15)

# Why W2 Wins
rect(s, 0.5, 3.7, 12.3, 2.5, LGREY)
txt(s, "Why W2 (ETA Inquiries) Was Chosen", 0.8, 4, 11.7, 0.4, size=20, bold=True, color=GREEN)
bullet_list(s, [
    "Highest volume (400/day) + lowest build risk = fastest win",
    "Simple: Lookup order + ask driver location via CRM API",
    "Saves 2.8 FTE (558 hours/month) - fastest payback",
    "Rebuilds trust after failed chatbot and RPA projects"
], 1, 4.6, 11.3, 1.5, size=15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: DECISION MATRIX
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "All 4 Workstreams Compared", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Compact table
headers = [
    ("Workstream", 2.8, BLUE),
    ("Volume", 1.5, BLUE),
    ("Savings", 1.5, BLUE),
    ("Risk", 1.3, BLUE),
    ("ROI", 1.6, BLUE),
    ("Decision", 2, BLUE)
]

x_pos = 0.5
y = 1.5
for header, width, color in headers:
    rect(s, x_pos, y, width, 0.5, color)
    txt(s, header, x_pos+0.1, y+0.1, width-0.2, 0.3, size=14, bold=True, color=WHITE)
    x_pos += width

# W2 Row (Winner)
y = 2
rect(s, 0.5, y, 2.8, 0.6, GREEN)
txt(s, "W2: ETA Inquiries", 0.6, y+0.15, 2.6, 0.3, size=14, bold=True, color=WHITE)
rect(s, 3.3, y, 1.5, 0.6, LGREY)
txt(s, "400/day", 3.4, y+0.15, 1.3, 0.3, size=14, color=DARK)
rect(s, 4.8, y, 1.5, 0.6, LGREY)
txt(s, "2.8 FTE", 4.9, y+0.15, 1.3, 0.3, size=14, color=DARK)
rect(s, 6.3, y, 1.3, 0.6, LGREY)
txt(s, "LOW", 6.4, y+0.15, 1.1, 0.3, size=14, bold=True, color=GREEN)
rect(s, 7.6, y, 1.6, 0.6, LGREY)
txt(s, "1-2 mo", 7.7, y+0.15, 1.4, 0.3, size=14, bold=True, color=GREEN)
rect(s, 9.2, y, 2, 0.6, GREEN)
txt(s, "PHASE 1", 9.3, y+0.15, 1.8, 0.3, size=14, bold=True, color=WHITE)

# W1 Row
y = 2.6
rect(s, 0.5, y, 2.8, 0.5, LGREY)
txt(s, "W1: Exceptions", 0.6, y+0.1, 2.6, 0.3, size=13, color=DARK)
rect(s, 3.3, y, 1.5, 0.5, LGREY)
txt(s, "180/day", 3.4, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 4.8, y, 1.5, 0.5, LGREY)
txt(s, "2.0 FTE", 4.9, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 6.3, y, 1.3, 0.5, LGREY)
txt(s, "MED", 6.4, y+0.1, 1.1, 0.3, size=13, color=AMBER)
rect(s, 7.6, y, 1.6, 0.5, LGREY)
txt(s, "6-12 mo", 7.7, y+0.1, 1.4, 0.3, size=13, color=DARK)
rect(s, 9.2, y, 2, 0.5, LGREY)
txt(s, "Defer", 9.3, y+0.1, 1.8, 0.3, size=13, color=MID)

# W4 Row
y = 3.1
rect(s, 0.5, y, 2.8, 0.5, LGREY)
txt(s, "W4: Billing", 0.6, y+0.1, 2.6, 0.3, size=13, color=DARK)
rect(s, 3.3, y, 1.5, 0.5, LGREY)
txt(s, "60/day", 3.4, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 4.8, y, 1.5, 0.5, LGREY)
txt(s, "1.5 FTE", 4.9, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 6.3, y, 1.3, 0.5, LGREY)
txt(s, "HIGH", 6.4, y+0.1, 1.1, 0.3, size=13, bold=True, color=RED)
rect(s, 7.6, y, 1.6, 0.5, LGREY)
txt(s, "3-6 mo", 7.7, y+0.1, 1.4, 0.3, size=13, color=DARK)
rect(s, 9.2, y, 2, 0.5, AMBER)
txt(s, "Phase 2", 9.3, y+0.1, 1.8, 0.3, size=13, bold=True, color=WHITE)

# W3 Row
y = 3.6
rect(s, 0.5, y, 2.8, 0.5, LGREY)
txt(s, "W3: Route Changes", 0.6, y+0.1, 2.6, 0.3, size=13, color=DARK)
rect(s, 3.3, y, 1.5, 0.5, LGREY)
txt(s, "90/day", 3.4, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 4.8, y, 1.5, 0.5, LGREY)
txt(s, "0.9 FTE", 4.9, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 6.3, y, 1.3, 0.5, LGREY)
txt(s, "HIGH", 6.4, y+0.1, 1.1, 0.3, size=13, bold=True, color=RED)
rect(s, 7.6, y, 1.6, 0.5, LGREY)
txt(s, "12-18 mo", 7.7, y+0.1, 1.4, 0.3, size=13, color=DARK)
rect(s, 9.2, y, 2, 0.5, LGREY)
txt(s, "Out of scope", 9.3, y+0.1, 1.8, 0.3, size=13, color=RED)

# Why not the others
txt(s, "Why not W1, W3, W4?", 0.8, 4.5, 11.7, 0.3, size=16, bold=True, color=NAVY)
bullet_list(s, [
    "W1: Dispatch console API limited - MEDIUM risk",
    "W4: Legacy Aurum billing has NO API, 24h lag - HIGH risk, defer to Phase 2",
    "W3: Route optimization complex, needs judgment - OUT OF SCOPE"
], 1, 4.9, 11.3, 1.5, size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: 3 CRITICAL QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, RED)
txt(s, "3 Critical Questions for Sarah", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

txt(s, "Answers will change the agent design", 0.5, 1.2, 12, 0.3, size=16, italic=True, color=MID, align=PP_ALIGN.CENTER)

# Q1
rect(s, 0.5, 1.8, 12.3, 1.3, LAMBER)
txt(s, "Q1: Driver App API - GPS Query or Message-Only?", 0.7, 2.0, 11.8, 0.3, size=15, bold=True, color=AMBER)
txt(s, "Can agent query GPS directly, or must it message driver and wait 2 min? If no API, agent is lookup-only.",
    0.7, 2.4, 11.8, 0.7, size=13, color=DARK)

# Q2
rect(s, 0.5, 3.3, 12.3, 1.3, LAMBER)
txt(s, "Q2: CRM Status Freshness - Real-time or Batch?", 0.7, 3.5, 11.8, 0.3, size=15, bold=True, color=AMBER)
txt(s, "How long after driver scan does status appear in Salesforce? If end-of-day batch, ETA logic breaks.",
    0.7, 3.9, 11.8, 0.7, size=13, color=DARK)

# Q3
rect(s, 0.5, 4.8, 12.3, 1.3, LAMBER)
txt(s, "Q3: Sandra Credit Audit Gap - Known Issue?", 0.7, 5.0, 11.8, 0.3, size=15, bold=True, color=AMBER)
txt(s, "£170 credit has no audit log. If exports are incomplete, Phase 2 billing agent will fail. Compliance risk.",
    0.7, 5.4, 11.8, 0.7, size=13, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: IMPLEMENTATION STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, GREEN)
txt(s, "Implementation Strategy", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Phase 1
rect(s, 1, 1.8, 11.3, 1.8, GREEN)
txt(s, "PHASE 1: Build ETA Inquiry Agent (W2)", 1.3, 2.0, 10.7, 0.4, size=20, bold=True, color=WHITE)
bullet_list(s, [
    "Quick operational win - rebuild trust after 2 failures",
    "400 inquiries/day automated - 2.8 FTE freed",
    "1-2 month ROI payback - prove value fast"
], 1.5, 2.5, 10, 1, size=16)

# Phase 2
rect(s, 1, 3.9, 11.3, 1.8, AMBER)
txt(s, "PHASE 2: Billing Disputes Agent (W4)", 1.3, 4.1, 10.7, 0.4, size=20, bold=True, color=WHITE)
bullet_list(s, [
    "Higher strategic value - compliance risk reduction",
    "Only after Phase 1 proves success - higher complexity",
    "Address legacy Aurum constraints with validated approach"
], 1.5, 4.6, 10, 1, size=16)

txt(s, "Bottom Line: Start with easy, high-volume win before tackling harder strategic problems",
    0.5, 6.2, 12.3, 0.5, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save("Gate2-Final-Presentation.pptx")
print("Merged presentation created: Gate2-Final-Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")

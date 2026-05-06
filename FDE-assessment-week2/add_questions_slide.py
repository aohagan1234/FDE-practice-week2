"""Add slide 5 with 3 critical discovery questions to Gate2-Presentation.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
ORANGE = RGBColor(0xE3, 0x6B, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
MID    = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xF2, 0xF4, 0xF7)
RED    = RGBColor(0xBF, 0x35, 0x2A)

# Load existing presentation
prs = Presentation('Gate2-Presentation.pptx')
BLANK = prs.slide_layouts[6]

# Helper functions
def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

# Move existing slides 5 and 6 to positions 6 and 7
# We'll insert the new slide at position 4 (index 4 = slide 5)
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)

# Insert new slide at position 4 (becomes slide 5)
s = prs.slides.add_slide(BLANK)

# Move it to position 4 (before current slide 5)
xml_slides.remove(slides[-1])  # Remove the new slide from end
xml_slides.insert(4, slides[-1])  # Insert at position 4

# Now get the slide at the correct position
s = prs.slides[4]

# ══════════════════════════════════════════════════════════════════════════════
# NEW SLIDE 5: 3 CRITICAL QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════

rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "3 Critical Questions for Sarah", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

txt(s, "Answers will materially change the agent design", 0.5, 1.2, 12, 0.3, size=16, italic=True, color=MID, align=PP_ALIGN.CENTER)

# Question 1
rect(s, 0.5, 1.8, 12.3, 1.4, LGREY)
txt(s, "Q1: Driver App API — GPS Query or Message-Only?", 0.7, 2.0, 11.8, 0.3, size=16, bold=True, color=NAVY)
txt(s, "Can the agent programmatically query driver GPS for real-time location, or must it send a message and wait for driver reply?",
    0.7, 2.4, 11.8, 0.5, size=13, color=DARK)
txt(s, "Impact: If message-only, ETA response time has a 2-min floor. If no API, agent becomes lookup-only tool.",
    0.7, 2.8, 11.8, 0.4, size=12, color=RED, italic=True)

# Question 2
rect(s, 0.5, 3.4, 12.3, 1.4, LGREY)
txt(s, "Q2: CRM Status Freshness — Real-time or Batch?", 0.7, 3.6, 11.8, 0.3, size=16, bold=True, color=NAVY)
txt(s, "When a driver scans a package, how long before that status appears in Salesforce? Immediate webhook, 5-30 min polling, or end-of-shift batch?",
    0.7, 4.0, 11.8, 0.5, size=13, color=DARK)
txt(s, "Impact: If end-of-day batch, 'out-for-delivery' status is unreliable. ETA logic breaks, requires full redesign.",
    0.7, 4.4, 11.8, 0.4, size=12, color=RED, italic=True)

# Question 3
rect(s, 0.5, 5.0, 12.3, 1.4, LGREY)
txt(s, "Q3: The Sandra Credit Audit Gap — Known Issue or Off-Policy?", 0.7, 5.2, 11.8, 0.3, size=16, bold=True, color=NAVY)
txt(s, "The 170 pound Hayes & Sons credit has no audit log entry. Is this a known gap where manual overrides bypass the export, or was it off-policy?",
    0.7, 5.6, 11.8, 0.5, size=13, color=DARK)
txt(s, "Impact: If exports are incomplete, Phase 2 billing agent will misidentify resolved disputes. Compliance risk.",
    0.7, 6.0, 11.8, 0.4, size=12, color=RED, italic=True)

txt(s, "These questions target assumptions rated LOW confidence - answers change architecture",
    0.5, 6.7, 12.3, 0.3, size=14, color=BLUE, align=PP_ALIGN.CENTER)

# Save
prs.save('Gate2-Presentation.pptx')
print("Slide 5 added: 3 Critical Questions for Sarah")
print("Total slides now: " + str(len(prs.slides)))

"""Expand Gate2 presentation with more workstream context and detail"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
ORANGE = RGBColor(0xE3, 0x6B, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
MID    = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xF2, 0xF4, 0xF7)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
RED    = RGBColor(0xBF, 0x35, 0x2A)
AMBER  = RGBColor(0xCC, 0x84, 0x00)
LAMBER = RGBColor(0xFD, 0xF0, 0xD5)
LRED   = RGBColor(0xF7, 0xDC, 0xDA)
LGREEN = RGBColor(0xD6, 0xEE, 0xDC)
TEAL   = RGBColor(0x0E, 0x7C, 0x7B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet_list(s, items, x, y, w, h, size=16, color=DARK):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.level = 0
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
rect(s, 0, 0, 13.33, 1.5, NAVY)
txt(s, "Apex Distribution Ltd", 0.5, 0.4, 12, 0.7, size=40, bold=True, color=WHITE)
txt(s, "Gate 2 Agentic Transformation Assessment", 0.5, 2.8, 12, 0.5, size=28, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(s, "Recommendation: Build ETA Inquiry Agent (Phase 1)", 0.5, 3.6, 12, 0.4, size=22, color=BLUE, align=PP_ALIGN.CENTER)
txt(s, "Then Tackle Billing Disputes (Phase 2)", 0.5, 4.2, 12, 0.4, size=20, color=AMBER, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "The Problem at Apex Distribution", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

txt(s, "Organization Context", 0.8, 1.5, 5, 0.3, size=18, bold=True, color=NAVY)
bullet_list(s, [
    "Birmingham-based carrier",
    "35-person Customer Ops team",
    "730 customer interactions/day",
    "4 work streams overlap constantly"
], 0.8, 1.9, 5, 1.5, size=15)

txt(s, "The Challenge", 7.3, 1.5, 5, 0.3, size=18, bold=True, color=RED)
bullet_list(s, [
    "Team overwhelmed with volume",
    "2 automation failures (chatbot, RPA)",
    "COO Sarah skeptical but needs results",
    "Legacy systems constrain options"
], 7.3, 1.9, 5, 1.5, size=15)

# Assessment approach
rect(s, 0.5, 3.7, 12.3, 2.5, LGREY)
txt(s, "Our Assessment Approach", 0.8, 4, 11.7, 0.4, size=20, bold=True, color=TEAL)
bullet_list(s, [
    "Analyzed all 4 work streams using ATX methodology (7-phase assessment)",
    "Scored each on: Volume × Automation Potential × Build Risk × Strategic Value",
    "Identified which work can be fully automated vs. requires human judgment",
    "Prioritized based on fastest ROI and lowest risk to rebuild trust"
], 1, 4.6, 11.3, 1.8, size=15)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE 4 WORKSTREAMS
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "4 Workstreams Analyzed", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# W1
rect(s, 0.5, 1.5, 6, 1.3, LGREY)
txt(s, "W1: Delivery Exceptions", 0.7, 1.7, 5.5, 0.3, size=17, bold=True, color=NAVY)
txt(s, "180/day | 12 min avg | 792 hrs/month", 0.7, 2.0, 5.5, 0.3, size=13, color=MID)
txt(s, "Driver issues, refused deliveries, damages, missed windows. Requires dispatcher judgment.",
    0.7, 2.3, 5.5, 0.5, size=12, color=DARK)

# W2 (highlighted)
rect(s, 6.8, 1.5, 6, 1.3, LGREEN)
txt(s, "W2: ETA Inquiries (WINNER)", 7, 1.7, 5.5, 0.3, size=17, bold=True, color=GREEN)
txt(s, "400/day | 4 min avg | 587 hrs/month", 7, 2.0, 5.5, 0.3, size=13, color=DARK)
txt(s, "'Where's my delivery?' Mostly lookup + driver location query. 95% automatable.",
    7, 2.3, 5.5, 0.5, size=12, color=DARK)

# W3
rect(s, 0.5, 3.1, 6, 1.3, LGREY)
txt(s, "W3: Dispatch Adjustments", 0.7, 3.3, 5.5, 0.3, size=17, bold=True, color=NAVY)
txt(s, "90/day | 18 min avg | 594 hrs/month", 0.7, 3.6, 5.5, 0.3, size=13, color=MID)
txt(s, "Mid-route changes, driver swaps, diversions. High judgment + route optimization complexity.",
    0.7, 3.9, 5.5, 0.5, size=12, color=DARK)

# W4
rect(s, 6.8, 3.1, 6, 1.3, LGREY)
txt(s, "W4: Billing Disputes", 7, 3.3, 5.5, 0.3, size=17, bold=True, color=NAVY)
txt(s, "60/day | 28 min avg | 616 hrs/month", 7, 3.6, 5.5, 0.3, size=13, color=MID)
txt(s, "Charge disputes, fuel surcharges, credits. High strategic value but legacy Aurum dependency.",
    7, 3.9, 5.5, 0.5, size=12, color=DARK)

# Summary
txt(s, "Total: 730 interactions/day consuming ~2,589 hours/month (74% of team capacity)",
    0.5, 4.8, 12.3, 0.4, size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

txt(s, "Analysis Question: Which workstream offers highest volume + lowest risk + fastest ROI?",
    0.5, 5.3, 12.3, 0.4, size=15, italic=True, color=MID, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: DECISION MATRIX
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "Decision Matrix: All 4 Workstreams Compared", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Table headers
headers = [
    ("Workstream", 2.8, BLUE),
    ("Volume", 1.5, BLUE),
    ("Savings", 1.5, BLUE),
    ("Risk", 1.3, BLUE),
    ("ROI", 1.6, BLUE),
    ("Decision", 2, BLUE)
]

x_pos = 0.5
y = 1.5
for header, width, color in headers:
    rect(s, x_pos, y, width, 0.5, color)
    txt(s, header, x_pos+0.1, y+0.1, width-0.2, 0.3, size=14, bold=True, color=WHITE)
    x_pos += width

# W2 Row (Winner)
y = 2
rect(s, 0.5, y, 2.8, 0.6, GREEN)
txt(s, "W2: ETA Inquiries", 0.6, y+0.15, 2.6, 0.3, size=14, bold=True, color=WHITE)
rect(s, 3.3, y, 1.5, 0.6, LGREY)
txt(s, "400/day", 3.4, y+0.15, 1.3, 0.3, size=14, color=DARK)
rect(s, 4.8, y, 1.5, 0.6, LGREY)
txt(s, "2.8 FTE", 4.9, y+0.15, 1.3, 0.3, size=14, color=DARK)
rect(s, 6.3, y, 1.3, 0.6, LGREY)
txt(s, "LOW", 6.4, y+0.15, 1.1, 0.3, size=14, bold=True, color=GREEN)
rect(s, 7.6, y, 1.6, 0.6, LGREY)
txt(s, "1-2 mo", 7.7, y+0.15, 1.4, 0.3, size=14, bold=True, color=GREEN)
rect(s, 9.2, y, 2, 0.6, GREEN)
txt(s, "PHASE 1", 9.3, y+0.15, 1.8, 0.3, size=14, bold=True, color=WHITE)

# W1 Row
y = 2.6
rect(s, 0.5, y, 2.8, 0.5, LGREY)
txt(s, "W1: Exceptions", 0.6, y+0.1, 2.6, 0.3, size=13, color=DARK)
rect(s, 3.3, y, 1.5, 0.5, LGREY)
txt(s, "180/day", 3.4, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 4.8, y, 1.5, 0.5, LGREY)
txt(s, "2.0 FTE", 4.9, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 6.3, y, 1.3, 0.5, LGREY)
txt(s, "MED", 6.4, y+0.1, 1.1, 0.3, size=13, color=AMBER)
rect(s, 7.6, y, 1.6, 0.5, LGREY)
txt(s, "6-12 mo", 7.7, y+0.1, 1.4, 0.3, size=13, color=DARK)
rect(s, 9.2, y, 2, 0.5, LGREY)
txt(s, "Defer", 9.3, y+0.1, 1.8, 0.3, size=13, color=MID)

# W4 Row
y = 3.1
rect(s, 0.5, y, 2.8, 0.5, LGREY)
txt(s, "W4: Billing", 0.6, y+0.1, 2.6, 0.3, size=13, color=DARK)
rect(s, 3.3, y, 1.5, 0.5, LGREY)
txt(s, "60/day", 3.4, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 4.8, y, 1.5, 0.5, LGREY)
txt(s, "1.5 FTE", 4.9, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 6.3, y, 1.3, 0.5, LGREY)
txt(s, "HIGH", 6.4, y+0.1, 1.1, 0.3, size=13, bold=True, color=RED)
rect(s, 7.6, y, 1.6, 0.5, LGREY)
txt(s, "3-6 mo", 7.7, y+0.1, 1.4, 0.3, size=13, color=DARK)
rect(s, 9.2, y, 2, 0.5, AMBER)
txt(s, "Phase 2", 9.3, y+0.1, 1.8, 0.3, size=13, bold=True, color=WHITE)

# W3 Row
y = 3.6
rect(s, 0.5, y, 2.8, 0.5, LGREY)
txt(s, "W3: Route Changes", 0.6, y+0.1, 2.6, 0.3, size=13, color=DARK)
rect(s, 3.3, y, 1.5, 0.5, LGREY)
txt(s, "90/day", 3.4, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 4.8, y, 1.5, 0.5, LGREY)
txt(s, "0.9 FTE", 4.9, y+0.1, 1.3, 0.3, size=13, color=DARK)
rect(s, 6.3, y, 1.3, 0.5, LGREY)
txt(s, "HIGH", 6.4, y+0.1, 1.1, 0.3, size=13, bold=True, color=RED)
rect(s, 7.6, y, 1.6, 0.5, LGREY)
txt(s, "12-18 mo", 7.7, y+0.1, 1.4, 0.3, size=13, color=DARK)
rect(s, 9.2, y, 2, 0.5, LGREY)
txt(s, "Out of scope", 9.3, y+0.1, 1.8, 0.3, size=13, color=RED)

txt(s, "W2 wins: Highest volume + Lowest risk = Fastest trust-building win after 2 prior failures",
    0.5, 4.4, 12.3, 0.4, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: WHY W2 WINS
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, GREEN)
txt(s, "Why W2 (ETA Inquiries) Is The Right Phase 1 Target", 0.5, 0.3, 12, 0.5, size=30, bold=True, color=WHITE)

# Left: Technical reasons
rect(s, 0.5, 1.5, 6, 4.5, LGREEN)
txt(s, "Technical Fit", 0.8, 1.8, 5.5, 0.3, size=18, bold=True, color=GREEN)
bullet_list(s, [
    "Highest volume: 400 inquiries/day",
    "95% automation potential - mostly lookup + driver query",
    "Simple integration: CRM REST API available",
    "No legacy dependencies (unlike Aurum billing)",
    "Clear escalation path: driver non-response in 2 min",
    "Saves 558 hours/month (2.8 FTE equivalent)"
], 0.8, 2.2, 5.3, 3.5, size=14)

# Right: Strategic reasons
rect(s, 6.8, 1.5, 6, 4.5, LGREEN)
txt(s, "Strategic Fit", 7.1, 1.8, 5.5, 0.3, size=18, bold=True, color=GREEN)
bullet_list(s, [
    "Fastest ROI: 1-2 month payback",
    "Rebuilds trust after 2 failed projects",
    "Proves agent value before tackling complex W4",
    "No chatbot interface - replies via existing channels",
    "Quick operational win validates approach",
    "Low risk = high confidence delivery"
], 7.1, 2.2, 5.3, 3.5, size=14)

txt(s, "After W2 success, use proven trust and capability to tackle W4 (billing disputes - higher value, higher complexity)",
    0.5, 6.3, 12.3, 0.6, size=15, italic=True, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: WHY NOT THE OTHERS - W1 & W3
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "Why Not W1 or W3 for Phase 1?", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# W1: Delivery Exceptions
rect(s, 0.5, 1.5, 12.3, 2.2, LAMBER)
txt(s, "W1: Delivery Exceptions - DEFERRED", 0.8, 1.7, 11.7, 0.4, size=20, bold=True, color=AMBER)
txt(s, "Volume: 180/day | Savings: 2.0 FTE | Risk: MEDIUM | ROI: 6-12 months",
    0.8, 2.1, 11.7, 0.3, size=14, color=DARK)

txt(s, "Why defer:", 0.8, 2.5, 11.7, 0.3, size=15, bold=True, color=AMBER)
bullet_list(s, [
    "Dispatch console has 'limited API surface' - unclear if write access exists",
    "Damage assessment requires human judgment (insurance, shipper coordination)",
    "SOP references retired system (DispatchHub) - documented process outdated",
    "Higher complexity than W2 with similar volume - not worth the risk for Phase 1"
], 1, 2.8, 11.3, 0.8, size=13)

# W3: Dispatch Adjustments
rect(s, 0.5, 4, 12.3, 2.2, LRED)
txt(s, "W3: Dispatch Adjustments - OUT OF SCOPE", 0.8, 4.2, 11.7, 0.4, size=20, bold=True, color=RED)
txt(s, "Volume: 90/day | Savings: 0.9 FTE | Risk: HIGH | ROI: 12-18 months",
    0.8, 4.6, 11.7, 0.3, size=14, color=DARK)

txt(s, "Why out of scope:", 0.8, 5.0, 11.7, 0.3, size=15, bold=True, color=RED)
bullet_list(s, [
    "Requires route optimization engine - unknown if exists or has API",
    "Tight time pressure + high judgment (driver swaps, diversions)",
    "Dispatch console API surface must be confirmed before scoping (Discovery Q6)",
    "If API is read-only or non-existent, automation limited to data surfacing only"
], 1, 5.3, 11.3, 0.8, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: WHY NOT W4 FOR PHASE 1
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, AMBER)
txt(s, "Why Not W4 (Billing Disputes) for Phase 1?", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

txt(s, "W4 has HIGHER strategic value but HIGHER risk - perfect for Phase 2 after W2 proves value",
    0.5, 1.2, 12.3, 0.4, size=16, italic=True, color=AMBER, align=PP_ALIGN.CENTER)

# The case FOR W4
rect(s, 0.5, 1.8, 6, 2, LGREEN)
txt(s, "Why W4 Is Important", 0.7, 2.0, 5.5, 0.3, size=16, bold=True, color=GREEN)
bullet_list(s, [
    "High strategic value: compliance risk reduction",
    "Audit trail gaps found (Sandra credit with no log)",
    "Customer satisfaction impact (28 min avg handling)",
    "Saves 1.5 FTE (308 hours/month)",
    "3-6 month ROI payback"
], 0.7, 2.4, 5.5, 1.3, size=13)

# The case AGAINST W4 for Phase 1
rect(s, 6.8, 1.8, 6, 2, LRED)
txt(s, "Why Defer to Phase 2", 7, 2.0, 5.5, 0.3, size=16, bold=True, color=RED)
bullet_list(s, [
    "Legacy Aurum billing: NO API, batch-only, 24h lag",
    "Prior RPA broke when Aurum schema changed",
    "Multiple async dependencies (shipper 24-72h, insurance)",
    "Lower volume than W2 (60/day vs 400/day)",
    "Need proven trust before tackling this complexity"
], 7, 2.4, 5.5, 1.3, size=13)

# The constraints
rect(s, 0.5, 4.1, 12.3, 2.2, LGREY)
txt(s, "The Aurum Constraint", 0.8, 4.3, 11.7, 0.3, size=18, bold=True, color=RED)
txt(s, "Aurum Billing System (on-prem Oracle, since 2008):",
    0.8, 4.7, 11.7, 0.3, size=14, color=DARK)
bullet_list(s, [
    "Batch-file exports only: Daily 02:00-04:00 GMT to CSV",
    "No real-time API - reconciliation lags 24 hours behind invoice generation",
    "Invoice modifications require manual ticket to Aurum support (48h turnaround)",
    "Schema changes quarterly without notice - broke previous RPA automation",
    "Any Phase 2 billing agent must include schema validation as first-class feature"
], 1, 5.0, 11.3, 1.2, size=13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: 3 CRITICAL QUESTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, RED)
txt(s, "3 Critical Questions for Sarah", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

txt(s, "Answers will materially change the agent design - rated LOW confidence",
    0.5, 1.2, 12, 0.3, size=16, italic=True, color=MID, align=PP_ALIGN.CENTER)

# Q1
rect(s, 0.5, 1.8, 12.3, 1.3, LAMBER)
txt(s, "Q1: Driver App API Surface - GPS Query or Message-Only?", 0.7, 2.0, 11.8, 0.3, size=15, bold=True, color=AMBER)
txt(s, "Can the agent programmatically query driver GPS for real-time location, or must it send a message and wait for driver reply (2 min timeout)?",
    0.7, 2.35, 11.8, 0.4, size=12, color=DARK)
txt(s, "Impact: If message-only, response time has 2-min floor. If no API at all, agent becomes lookup-only tool (major capability reduction).",
    0.7, 2.75, 11.8, 0.35, size=11, italic=True, color=RED)

# Q2
rect(s, 0.5, 3.3, 12.3, 1.3, LAMBER)
txt(s, "Q2: CRM Status Freshness - Real-time Webhook or Batch?", 0.7, 3.5, 11.8, 0.3, size=15, bold=True, color=AMBER)
txt(s, "When a driver scans a package on delivery, how long before that status appears in Salesforce CRM? Immediate (<1 min), polling (5-30 min), or end-of-shift batch?",
    0.7, 3.85, 11.8, 0.4, size=12, color=DARK)
txt(s, "Impact: If end-of-day batch, 'out-for-delivery' status is unreliable. Agent cannot distinguish 'in transit' from 'already delivered'. ETA logic breaks - requires full redesign.",
    0.7, 4.25, 11.8, 0.35, size=11, italic=True, color=RED)

# Q3
rect(s, 0.5, 4.8, 12.3, 1.3, LAMBER)
txt(s, "Q3: Sandra Credit Audit Trail Gap - Known Issue or Off-Policy?", 0.7, 5.0, 11.8, 0.3, size=15, bold=True, color=AMBER)
txt(s, "The £170 Hayes & Sons goodwill credit has no audit log entry. Is this a known gap where manual overrides bypass the export, or was it off-policy?",
    0.7, 5.35, 11.8, 0.4, size=12, color=DARK)
txt(s, "Impact: If exports are incomplete, Phase 2 billing agent will misidentify resolved disputes as open. Compliance exposure - credits without audit trail.",
    0.7, 5.75, 11.8, 0.35, size=11, italic=True, color=RED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: IMPLEMENTATION STRATEGY
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 1, TEAL)
txt(s, "Phased Implementation Strategy", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Phase 1
rect(s, 1, 1.5, 11.3, 2, GREEN)
txt(s, "PHASE 1: ETA Inquiry Agent (W2) - 6-8 Week Build", 1.3, 1.7, 10.7, 0.4, size=20, bold=True, color=WHITE)
bullet_list(s, [
    "Scope: Automate 'where's my delivery' inquiries (400/day)",
    "Quick win: Rebuild Sarah's trust after chatbot and RPA failures",
    "Low risk: CRM REST API available, no legacy dependencies",
    "1-2 month ROI payback: 2.8 FTE savings (558 hours/month)",
    "Success metrics: <5 min response, >80% resolution without escalation"
], 1.5, 2.2, 10, 1.2, size=14)

# Phase 2
rect(s, 1, 3.8, 11.3, 2, AMBER)
txt(s, "PHASE 2: Billing Dispute Agent (W4) - After Phase 1 Success", 1.3, 4.0, 10.7, 0.4, size=20, bold=True, color=WHITE)
bullet_list(s, [
    "Scope: Triage, investigate, and resolve billing disputes (60/day)",
    "Strategic value: Compliance risk reduction + audit trail enforcement",
    "Higher complexity: Aurum batch dependency + schema validation required",
    "3-6 month ROI payback: 1.5 FTE savings + penalty avoidance",
    "Only proceed after Phase 1 proves agent reliability and team buy-in"
], 1.5, 4.5, 10, 1.2, size=14)

txt(s, "Bottom Line: Start easy, prove value, then tackle strategic complexity",
    0.5, 6.2, 12.3, 0.5, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save("Gate2-Final-Presentation.pptx")
print("Expanded presentation created: Gate2-Final-Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")

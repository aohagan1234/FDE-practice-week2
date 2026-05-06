"""Generate high-level summary presentation — Apex Distribution Ltd."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
BLUE   = RGBColor(0x2E, 0x5F, 0xA3)
ORANGE = RGBColor(0xE3, 0x6B, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
MID    = RGBColor(0x60, 0x60, 0x60)
LGREY  = RGBColor(0xF2, 0xF4, 0xF7)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
RED    = RGBColor(0xBF, 0x35, 0x2A)
AMBER  = RGBColor(0xCC, 0x84, 0x00)
LAMBER = RGBColor(0xFD, 0xF0, 0xD5)
LRED   = RGBColor(0xF7, 0xDC, 0xDA)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def txt(s, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet_list(s, items, x, y, w, h, size=16, color=DARK):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]

        p.level = 0
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color

    return tb


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 7.5, LGREY)
rect(s, 0, 0, 13.33, 1.5, NAVY)

txt(s, "Apex Distribution Ltd", 0.5, 0.4, 12, 0.7, size=40, bold=True, color=WHITE)
txt(s, "Agentic Transformation Assessment — Gate 2", 0.5, 3, 12, 0.5, size=28, bold=True, color=NAVY)
txt(s, "Recommendation: Build ETA Inquiry Agent (W2) First", 0.5, 4, 12, 0.5, size=22, color=BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "The Problem", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

rect(s, 0.5, 1.5, 12.3, 5, WHITE)
txt(s, "Apex Distribution Ltd — Birmingham, UK", 0.8, 1.8, 11, 0.4, size=20, bold=True, color=NAVY)

bullet_list(s, [
    "35-person Customer Operations team overwhelmed",
    "Handling 730 customer requests daily",
    "Previous automation attempts failed (2024 chatbot, RPA project)",
    "COO Sarah Whitmore is skeptical but needs results"
], 1, 2.5, 11, 3, size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE 4 OPTIONS
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "4 Workstreams Analyzed", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# W1 Box
rect(s, 0.5, 1.5, 6, 1.2, LGREY)
txt(s, "W1: Delivery Exceptions", 0.7, 1.7, 5.5, 0.3, size=18, bold=True, color=NAVY)
txt(s, "180/day | Damaged packages, refusals", 0.7, 2.0, 5.5, 0.5, size=14, color=MID)

# W2 Box (highlighted)
rect(s, 7, 1.5, 6, 1.2, GREEN)
txt(s, "W2: ETA Inquiries ✓", 7.2, 1.7, 5.5, 0.3, size=18, bold=True, color=WHITE)
txt(s, "400/day | 'Where's my delivery?'", 7.2, 2.0, 5.5, 0.5, size=14, color=WHITE)

# W3 Box
rect(s, 0.5, 3, 6, 1.2, LGREY)
txt(s, "W3: Route Changes", 0.7, 3.2, 5.5, 0.3, size=18, bold=True, color=NAVY)
txt(s, "90/day | Mid-delivery adjustments", 0.7, 3.5, 5.5, 0.5, size=14, color=MID)

# W4 Box
rect(s, 7, 3, 6, 1.2, LGREY)
txt(s, "W4: Billing Disputes", 7.2, 3.2, 5.5, 0.3, size=18, bold=True, color=NAVY)
txt(s, "60/day | Customer charge disputes", 7.2, 3.5, 5.5, 0.5, size=14, color=MID)

txt(s, "Analysis Criteria: Volume × Automation Potential × Build Risk × ROI", 0.5, 4.8, 12, 0.4, size=16, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: WHY W2 WON
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 1, GREEN)
txt(s, "Why W2 (ETA Inquiries) Was Chosen", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

txt(s, "Highest Volume + Lowest Risk = Fastest Win", 0.8, 1.5, 11, 0.5, size=22, bold=True, color=GREEN)

bullet_list(s, [
    "400 requests/day — Biggest time drain (558 hours/month)",
    "Simple to build — Lookup order + ask driver for location",
    "Works with existing systems — Modern CRM has API available",
    "Fastest ROI — Pays back in 1-2 months",
    "Saves 2.8 FTE worth of work — Team can redeploy capacity"
], 1, 2.3, 11, 3.5, size=18)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: WHY NOT THE OTHERS
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "Why Not W1, W3, or W4?", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# W1
rect(s, 0.5, 1.5, 12.3, 1, LAMBER)
txt(s, "W1: Delivery Exceptions", 0.8, 1.7, 5, 0.3, size=16, bold=True, color=AMBER)
txt(s, "• Dispatch system has limited API access — build risk MEDIUM", 0.8, 2, 11, 0.4, size=14, color=DARK)

# W3
rect(s, 0.5, 2.7, 12.3, 1, LRED)
txt(s, "W3: Route Changes", 0.8, 2.9, 5, 0.3, size=16, bold=True, color=RED)
txt(s, "• Route optimization too complex, needs human judgment — build risk HIGH", 0.8, 3.2, 11, 0.4, size=14, color=DARK)

# W4
rect(s, 0.5, 3.9, 12.3, 1.3, LRED)
txt(s, "W4: Billing Disputes", 0.8, 4.1, 5, 0.3, size=16, bold=True, color=RED)
txt(s, "• Old billing system (Aurum) has NO API, 24-hour data lag, breaks when updated", 0.8, 4.4, 11, 0.4, size=14, color=DARK)
txt(s, "• Higher strategic value (compliance) but too risky for Phase 1", 0.8, 4.8, 11, 0.4, size=14, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: COMPARISON TABLE
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "Decision Matrix", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Table headers
rect(s, 0.5, 1.5, 2.5, 0.5, BLUE)
txt(s, "Workstream", 0.6, 1.6, 2.3, 0.3, size=14, bold=True, color=WHITE)

rect(s, 3, 1.5, 1.8, 0.5, BLUE)
txt(s, "Volume/Day", 3.1, 1.6, 1.6, 0.3, size=14, bold=True, color=WHITE)

rect(s, 4.8, 1.5, 1.8, 0.5, BLUE)
txt(s, "Time Savings", 4.9, 1.6, 1.6, 0.3, size=14, bold=True, color=WHITE)

rect(s, 6.6, 1.5, 1.8, 0.5, BLUE)
txt(s, "Build Risk", 6.7, 1.6, 1.6, 0.3, size=14, bold=True, color=WHITE)

rect(s, 8.4, 1.5, 2, 0.5, BLUE)
txt(s, "ROI Payback", 8.5, 1.6, 1.8, 0.3, size=14, bold=True, color=WHITE)

rect(s, 10.4, 1.5, 2.4, 0.5, BLUE)
txt(s, "Decision", 10.5, 1.6, 2.2, 0.3, size=14, bold=True, color=WHITE)

# W2 Row (Winner)
y = 2
rect(s, 0.5, y, 2.5, 0.6, GREEN)
txt(s, "W2: ETA Inquiries", 0.6, y+0.1, 2.3, 0.3, size=14, bold=True, color=WHITE)
rect(s, 3, y, 1.8, 0.6, LGREY)
txt(s, "400", 3.1, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 4.8, y, 1.8, 0.6, LGREY)
txt(s, "2.8 FTE", 4.9, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 6.6, y, 1.8, 0.6, LGREY)
txt(s, "LOW", 6.7, y+0.15, 1.6, 0.3, size=14, bold=True, color=GREEN)
rect(s, 8.4, y, 2, 0.6, LGREY)
txt(s, "1-2 months", 8.5, y+0.15, 1.8, 0.3, size=14, bold=True, color=GREEN)
rect(s, 10.4, y, 2.4, 0.6, GREEN)
txt(s, "✓ PHASE 1", 10.5, y+0.15, 2.2, 0.3, size=14, bold=True, color=WHITE)

# W1 Row
y = 2.6
rect(s, 0.5, y, 2.5, 0.6, LGREY)
txt(s, "W1: Exceptions", 0.6, y+0.1, 2.3, 0.3, size=14, color=DARK)
rect(s, 3, y, 1.8, 0.6, LGREY)
txt(s, "180", 3.1, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 4.8, y, 1.8, 0.6, LGREY)
txt(s, "2.0 FTE", 4.9, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 6.6, y, 1.8, 0.6, LGREY)
txt(s, "MEDIUM", 6.7, y+0.15, 1.6, 0.3, size=14, color=AMBER)
rect(s, 8.4, y, 2, 0.6, LGREY)
txt(s, "6-12 months", 8.5, y+0.15, 1.8, 0.3, size=14, color=DARK)
rect(s, 10.4, y, 2.4, 0.6, LGREY)
txt(s, "Defer", 10.5, y+0.15, 2.2, 0.3, size=14, color=MID)

# W4 Row
y = 3.2
rect(s, 0.5, y, 2.5, 0.6, LGREY)
txt(s, "W4: Billing", 0.6, y+0.1, 2.3, 0.3, size=14, color=DARK)
rect(s, 3, y, 1.8, 0.6, LGREY)
txt(s, "60", 3.1, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 4.8, y, 1.8, 0.6, LGREY)
txt(s, "1.5 FTE", 4.9, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 6.6, y, 1.8, 0.6, LGREY)
txt(s, "HIGH", 6.7, y+0.15, 1.6, 0.3, size=14, bold=True, color=RED)
rect(s, 8.4, y, 2, 0.6, LGREY)
txt(s, "3-6 months", 8.5, y+0.15, 1.8, 0.3, size=14, color=DARK)
rect(s, 10.4, y, 2.4, 0.6, AMBER)
txt(s, "Phase 2", 10.5, y+0.15, 2.2, 0.3, size=14, bold=True, color=WHITE)

# W3 Row
y = 3.8
rect(s, 0.5, y, 2.5, 0.6, LGREY)
txt(s, "W3: Route Changes", 0.6, y+0.1, 2.3, 0.3, size=14, color=DARK)
rect(s, 3, y, 1.8, 0.6, LGREY)
txt(s, "90", 3.1, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 4.8, y, 1.8, 0.6, LGREY)
txt(s, "0.9 FTE", 4.9, y+0.15, 1.6, 0.3, size=14, color=DARK)
rect(s, 6.6, y, 1.8, 0.6, LGREY)
txt(s, "HIGH", 6.7, y+0.15, 1.6, 0.3, size=14, bold=True, color=RED)
rect(s, 8.4, y, 2, 0.6, LGREY)
txt(s, "12-18 months", 8.5, y+0.15, 1.8, 0.3, size=14, color=DARK)
rect(s, 10.4, y, 2.4, 0.6, LGREY)
txt(s, "Out of scope", 10.5, y+0.15, 2.2, 0.3, size=14, color=RED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: THE STRATEGY
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 1, NAVY)
txt(s, "Implementation Strategy", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Phase 1
rect(s, 0.5, 1.5, 12.3, 1.5, GREEN)
txt(s, "Phase 1: Build ETA Inquiry Agent (W2)", 0.8, 1.7, 11, 0.4, size=20, bold=True, color=WHITE)
bullet_list(s, [
    "Quick win to rebuild trust after 2 failed automation attempts",
    "Prove the technology works with lowest-risk, highest-volume use case",
    "1-2 month payback — fast ROI demonstration"
], 1, 2.2, 11, 1, size=16)

# Phase 2
rect(s, 0.5, 3.3, 12.3, 1.5, AMBER)
txt(s, "Phase 2: Tackle Billing Disputes (W4)", 0.8, 3.5, 11, 0.4, size=20, bold=True, color=WHITE)
bullet_list(s, [
    "More strategic (compliance risk reduction)",
    "Higher complexity — requires proven trust from Phase 1 success",
    "Address legacy system constraints with validated approach"
], 1, 4, 11, 1, size=16)

txt(s, "Bottom Line: Start with the easy, high-volume win before attempting harder, strategic problems",
    0.5, 5.3, 12.3, 0.8, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: KEY BENEFITS
# ══════════════════════════════════════════════════════════════════════════════

s = new_slide()
rect(s, 0, 0, 13.33, 1, GREEN)
txt(s, "Expected Benefits — W2 ETA Agent", 0.5, 0.3, 12, 0.5, size=32, bold=True, color=WHITE)

# Benefit boxes
rect(s, 1, 2, 3.5, 1.5, LGREY)
txt(s, "Time Savings", 1.2, 2.2, 3, 0.3, size=18, bold=True, color=GREEN)
txt(s, "558 hours/month\n2.8 FTE equivalent", 1.2, 2.6, 3, 0.8, size=16, color=DARK, align=PP_ALIGN.CENTER)

rect(s, 4.9, 2, 3.5, 1.5, LGREY)
txt(s, "Speed", 4.9, 2.2, 3.5, 0.3, size=18, bold=True, color=GREEN)
txt(s, "<5 min response\nvs. 4 min current avg", 5.1, 2.6, 3, 0.8, size=16, color=DARK, align=PP_ALIGN.CENTER)

rect(s, 8.8, 2, 3.5, 1.5, LGREY)
txt(s, "Customer Satisfaction", 8.8, 2.2, 3.5, 0.3, size=18, bold=True, color=GREEN)
txt(s, "Faster, more accurate\nETA information", 9, 2.6, 3, 0.8, size=16, color=DARK, align=PP_ALIGN.CENTER)

rect(s, 2.75, 4, 3.5, 1.5, LGREY)
txt(s, "Scalability", 2.95, 4.2, 3, 0.3, size=18, bold=True, color=GREEN)
txt(s, "Handle volume spikes\nwithout hiring", 2.95, 4.6, 3, 0.8, size=16, color=DARK, align=PP_ALIGN.CENTER)

rect(s, 6.7, 4, 3.5, 1.5, LGREY)
txt(s, "Team Morale", 6.9, 4.2, 3, 0.3, size=18, bold=True, color=GREEN)
txt(s, "Staff focus on complex\nwork, not repetitive lookups", 6.9, 4.6, 3, 0.8, size=16, color=DARK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════

prs.save("Apex-Summary-Presentation.pptx")
print("Presentation saved: Apex-Summary-Presentation.pptx")
